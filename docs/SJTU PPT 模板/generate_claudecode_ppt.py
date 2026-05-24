#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SJTU Premium PPT Generator — Claude Code 源码解析
Wine Red + Gold 主题 · 模板布局 + 精品组件系统
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# ═══════════════════════════════════════════════════════════
# COLOR SYSTEM  —  Wine Red + Gold
# ═══════════════════════════════════════════════════════════
SJTU_RED      = RGBColor(0xA6, 0x20, 0x38)
SJTU_DEEP     = RGBColor(0x6B, 0x15, 0x25)
SJTU_ACCENT   = RGBColor(0xC0, 0x00, 0x00)

GOLD          = RGBColor(0xC5, 0xA4, 0x6C)
GOLD_LIGHT    = RGBColor(0xD4, 0xBB, 0x8A)
GOLD_PALE     = RGBColor(0xF0, 0xE6, 0xD2)

WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE     = RGBColor(0xFA, 0xFA, 0xFB)
WARM_GRAY_1   = RGBColor(0xF5, 0xF4, 0xF2)
WARM_GRAY_2   = RGBColor(0xED, 0xEC, 0xEA)
WARM_GRAY_3   = RGBColor(0xD2, 0xD0, 0xCD)

TEXT_PRIMARY   = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_BODY      = RGBColor(0x3D, 0x3D, 0x4E)
TEXT_CAPTION   = RGBColor(0x6B, 0x72, 0x80)
TEXT_SUBTLE    = RGBColor(0x9C, 0xA0, 0xAB)

BLUE_TEXT      = RGBColor(0x00, 0x40, 0x99)
SUCCESS        = RGBColor(0x2D, 0x8A, 0x56)
WARN_AMBER     = RGBColor(0xC4, 0x7F, 0x17)

TBL_HEADER     = SJTU_RED
TBL_ALT        = RGBColor(0xF2, 0xF1, 0xEE)

# ═══════════════════════════════════════════════════════════
# TYPE SCALE
# ═══════════════════════════════════════════════════════════
SZ_T = 24; SZ_H = 16; SZ_B = 13; SZ_L = 11; SZ_M = 8

FN = '微软雅黑'
LQ = '\u201c'; RQ = '\u201d'
def Q(t): return LQ + t + RQ

# ═══════════════════════════════════════════════════════════
# DIMENSIONS
# ═══════════════════════════════════════════════════════════
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.9)
CW      = SLIDE_W - Inches(1.8)

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'SJTU 模板.pptx')

# ═══════════════════════════════════════════════════════════
# LOAD TEMPLATE
# ═══════════════════════════════════════════════════════════
prs = Presentation(TEMPLATE_PATH)

while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn('r:id'))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

LAY_COVER   = prs.slide_layouts[0]
LAY_CONTENT = prs.slide_layouts[7]
LAY_END     = prs.slide_layouts[12]

# ═══════════════════════════════════════════════════════════
# SHAPE PRIMITIVES
# ═══════════════════════════════════════════════════════════

def _set_ea(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FN)

def add_content_slide():
    return prs.slides.add_slide(LAY_CONTENT)

def rect(slide, l, t, w, h, fill, border=None, bw=Pt(0.75)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = bw
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def rrect(slide, l, t, w, h, fill, border=None, bw=Pt(0.75), radius=0.03):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = bw
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    if radius is not None: sh.adjustments[0] = radius
    return sh

def oval(slide, l, t, w, h, fill, border=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def _make_oxml(tag):
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    prefix = tag.split(':')[0] if ':' in tag else None
    local = tag.split(':')[1] if ':' in tag else tag
    if prefix == 'a':
        return etree.Element(f'{{{nsmap["a"]}}}{local}')
    return etree.Element(tag)

# ═══════════════════════════════════════════════════════════
# TEXT HELPERS
# ═══════════════════════════════════════════════════════════

def stxt(tf, text, sz=SZ_B, color=TEXT_BODY, bold=False, al=PP_ALIGN.LEFT):
    tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = FN; p.alignment = al
    _set_ea(p.runs[0]) if p.runs else None
    return p

def apara(tf, text, sz=SZ_B, color=TEXT_BODY, bold=False, al=PP_ALIGN.LEFT, sb=Pt(4), sa=Pt(2)):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = FN
    p.alignment = al; p.space_before = sb; p.space_after = sa
    _set_ea(p.runs[0]) if p.runs else None
    return p

def bullets(tf, items, sz=SZ_B, color=TEXT_BODY):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
            for r_el in p._p.findall(qn('a:r')):
                p._p.remove(r_el)
        else:
            p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = '  \u25aa   '; r1.font.size = Pt(sz); r1.font.color.rgb = SJTU_RED; r1.font.name = FN
        _set_ea(r1)
        r2 = p.add_run()
        r2.text = item; r2.font.size = Pt(sz); r2.font.color.rgb = color; r2.font.name = FN
        _set_ea(r2)
        p.space_before = Pt(6); p.space_after = Pt(3); p.line_spacing = 1.3

# ═══════════════════════════════════════════════════════════
# PREMIUM COMPONENTS
# ═══════════════════════════════════════════════════════════

def add_shadow(shape):
    spPr = shape._element.spPr
    effl = spPr.find(qn('a:effectLst'))
    if effl is None:
        effl = _make_oxml('a:effectLst'); spPr.append(effl)
    else:
        for ch in list(effl): effl.remove(ch)
    os_el = _make_oxml('a:outerShdw')
    os_el.set('blurRad', '63500'); os_el.set('dist', '19050')
    os_el.set('dir', '5400000'); os_el.set('rotWithShape', '0')
    sc = _make_oxml('a:srgbClr'); sc.set('val', '6B1525')
    al = _make_oxml('a:alpha'); al.set('val', '12000')
    sc.append(al); os_el.append(sc); effl.append(os_el)

def sjtu_header(slide, title, page=None):
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 11:
            ph.text = ''
            run = ph.text_frame.paragraphs[0].add_run()
            run.text = title; run.font.size = Pt(22); run.font.bold = True
            run.font.color.rgb = WHITE; run.font.name = FN; _set_ea(run)
        elif idx == 12:
            ph.text = ''
    sjtu_footer(slide, page)

def sjtu_footer(slide, page=None):
    if page:
        y_ft = SLIDE_H - Inches(0.42)
        t = tb(slide, SLIDE_W - Inches(1.1), y_ft + Pt(2), Inches(0.7), Inches(0.25))
        tf = t.text_frame; p = tf.paragraphs[0]; p.text = ''
        r1 = p.add_run()
        r1.text = '\u25c6  '; r1.font.size = Pt(5); r1.font.color.rgb = SJTU_RED; r1.font.name = FN; _set_ea(r1)
        r2 = p.add_run()
        r2.text = f'{page:02d}'; r2.font.size = Pt(SZ_L); r2.font.color.rgb = TEXT_CAPTION; r2.font.name = FN; _set_ea(r2)
        p.alignment = PP_ALIGN.RIGHT

def card(slide, l, t, w, h, fill=WHITE, border=WARM_GRAY_2, gold_top=False, shadow=True):
    sh = rrect(slide, l, t, w, h, fill, border=border, bw=Pt(0.5))
    if shadow: add_shadow(sh)
    if gold_top: rect(slide, l, t, w, Pt(2.5), SJTU_RED)
    return sh

def make_tbl(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table

def scell(cell, text, sz=SZ_B, bold=False, color=TEXT_BODY, bg=None, al=PP_ALIGN.LEFT):
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = FN; p.alignment = al
    if p.runs: _set_ea(p.runs[0])
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg: cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.margin_left = Pt(8); cell.margin_right = Pt(8)
    cell.margin_top = Pt(5); cell.margin_bottom = Pt(5)

def tbl_hdr(table, headers):
    for i, h in enumerate(headers):
        scell(table.cell(0, i), h, sz=SZ_B, bold=True, color=WHITE, bg=TBL_HEADER, al=PP_ALIGN.CENTER)

def tbl_data(table, data, start=1):
    for ri, row in enumerate(data):
        rn = start + ri
        bg = TBL_ALT if rn % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            is_first = (ci == 0)
            scell(table.cell(rn, ci), val, sz=SZ_B, bold=is_first,
                  color=TEXT_PRIMARY if is_first else TEXT_BODY, bg=bg)

def _cb(cell, top=None, bottom=None, left=None, right=None):
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = _make_oxml('a:tcPr'); tc.append(tcPr)
    for tag, spec in [('lnT', top), ('lnB', bottom), ('lnL', left), ('lnR', right)]:
        for el in tcPr.findall(qn(f'a:{tag}')):
            tcPr.remove(el)
        ln = _make_oxml(f'a:{tag}')
        if spec is None:
            nf = _make_oxml('a:noFill'); ln.append(nf)
        else:
            hx, w = spec
            ln.set('w', str(w))
            sf = _make_oxml('a:solidFill'); sc = _make_oxml('a:srgbClr'); sc.set('val', hx)
            sf.append(sc); ln.append(sf)
        tcPr.append(ln)

def style_tbl(table, rows, cols):
    tbl = table._tbl; tblPr = tbl.tblPr
    for ch in list(tblPr):
        if 'tblStyle' in ch.tag: tblPr.remove(ch)
    tblPr.attrib['bandRow'] = '0'; tblPr.attrib['firstRow'] = '0'
    tblPr.attrib['lastRow'] = '0'; tblPr.attrib['bandCol'] = '0'
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            if r == 0:
                _cb(cell, top=('A62038', 6350), bottom=('A62038', 12700))
            elif r < rows - 1:
                _cb(cell, bottom=('EDECEA', 6350))
            else:
                _cb(cell)

def insight(slide, left, top, width, text, accent=GOLD):
    h = Inches(0.48)
    rrect(slide, left, top, width, h, WARM_GRAY_1, border=WARM_GRAY_2)
    rect(slide, left, top, Pt(3), h, accent)
    t = tb(slide, left + Pt(14), top, width - Pt(18), h)
    tf = t.text_frame; tf.word_wrap = True
    tf.margin_top = 0; tf.margin_bottom = 0; tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    p = tf.paragraphs[0]; p.text = ''; p.space_before = Pt(0); p.space_after = Pt(0)
    r1 = p.add_run()
    r1.text = '\u25b8  '; r1.font.size = Pt(SZ_L); r1.font.color.rgb = accent; r1.font.name = FN; _set_ea(r1)
    r2 = p.add_run()
    r2.text = text; r2.font.size = Pt(SZ_L); r2.font.color.rgb = TEXT_BODY; r2.font.name = FN; _set_ea(r2)
    try:
        bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None: bodyPr.set('anchor', 'ctr')
    except: pass

def highlight_card(slide, left, top, width, label, text):
    h = Inches(1.1)
    c = card(slide, left, top, width, h, WARM_GRAY_1, border=WARM_GRAY_2)
    rect(slide, left, top, Pt(4), h, GOLD)
    pill(slide, left + Inches(0.25), top + Inches(0.12), Inches(2.2), Inches(0.3), label, bg=GOLD, fg=SJTU_DEEP)
    t = tb(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.5), Inches(0.45))
    tf = t.text_frame; tf.word_wrap = True
    stxt(tf, text, sz=SZ_B, color=TEXT_PRIMARY, bold=True)
    return c

def gold_div(slide, left, top, width):
    mid = left + width // 2
    rect(slide, left, top, width, Pt(0.5), WARM_GRAY_3)
    d = Pt(6)
    sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, mid - d//2, top - d//2, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = SJTU_RED; sh.line.fill.background()

def pill(slide, l, t, w, h, text, bg=SJTU_RED, fg=WHITE):
    sh = rrect(slide, l, t, w, h, bg)
    stxt(sh.text_frame, text, sz=SZ_L, color=fg, bold=True, al=PP_ALIGN.CENTER)

def badge(slide, l, t, text, size=Inches(0.5), bg=SJTU_RED, fg=WHITE, sz=SZ_L):
    sh = oval(slide, l, t, size, size, bg)
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = fg
    p.font.bold = True; p.font.name = FN; p.alignment = PP_ALIGN.CENTER
    return sh

def stat_callout(slide, l, t, w, h, number, unit, label, desc=None, accent=SJTU_RED, num_sz=36):
    c = card(slide, l, t, w, h, WHITE, WARM_GRAY_2)
    rect(slide, l, t, w, Pt(3), accent)
    tlb = tb(slide, l + Inches(0.2), t + Inches(0.1), w - Inches(0.4), Inches(0.22))
    stxt(tlb.text_frame, label, sz=SZ_L, color=TEXT_CAPTION)
    tn = tb(slide, l + Inches(0.15), t + Inches(0.35), w - Inches(0.3), Inches(0.55))
    tf = tn.text_frame; p = tf.paragraphs[0]; p.text = ''
    r1 = p.add_run()
    r1.text = number; r1.font.size = Pt(num_sz); r1.font.color.rgb = accent
    r1.font.bold = True; r1.font.name = FN; _set_ea(r1)
    if unit:
        r2 = p.add_run()
        r2.text = ' ' + unit; r2.font.size = Pt(SZ_B); r2.font.color.rgb = TEXT_CAPTION
        r2.font.name = FN; _set_ea(r2)
    if desc:
        td = tb(slide, l + Inches(0.2), t + h - Inches(0.28), w - Inches(0.4), Inches(0.22))
        stxt(td.text_frame, desc, sz=SZ_M, color=TEXT_CAPTION)
    return c

def step_flow(slide, left, top, width, steps):
    n = len(steps); gap = int(width / n)
    for i, (num, label) in enumerate(steps):
        cx = int(left + gap * i + gap / 2) - Inches(0.18)
        cy = top
        if i > 0:
            prev_cx = int(left + gap * (i - 1) + gap / 2) + Inches(0.18)
            rect(slide, prev_cx, cy + Pt(12), cx - prev_cx, Pt(2), SJTU_RED)
        badge(slide, cx, cy, num, size=Inches(0.36), sz=SZ_L)
        t = tb(slide, cx - Inches(0.8), cy + Inches(0.42), Inches(1.96), Inches(0.3))
        stxt(t.text_frame, label, sz=SZ_L, color=TEXT_PRIMARY, al=PP_ALIGN.CENTER, bold=True)

def quote_block(slide, left, top, width, text, source=None):
    h = Inches(1.5)
    c = card(slide, left, top, width, h, WARM_GRAY_1, border=WARM_GRAY_2)
    rect(slide, left, top, Pt(4), h, GOLD)
    t1 = tb(slide, left + Inches(0.15), top - Inches(0.1), Inches(0.8), Inches(0.8))
    stxt(t1.text_frame, LQ, sz=60, color=GOLD, bold=True)
    t2 = tb(slide, left + Inches(0.6), top + Inches(0.25), width - Inches(1.2), Inches(0.7))
    t2.text_frame.word_wrap = True
    stxt(t2.text_frame, text, sz=SZ_H, color=TEXT_PRIMARY)
    if source:
        t3 = tb(slide, left + Inches(0.6), top + Inches(1.0), width - Inches(1.2), Inches(0.3))
        stxt(t3.text_frame, '\u2014\u2014 ' + source, sz=SZ_L, color=TEXT_CAPTION, al=PP_ALIGN.RIGHT)
    return c

def chevron_flow(slide, left, top, width, items):
    n = len(items); cw = int(width / n) - Pt(4)
    colors = [SJTU_RED, GOLD, RGBColor(0x8B, 0x5C, 0x3E), TEXT_PRIMARY]
    for i, (label, sub) in enumerate(items):
        x = int(left + int(width / n) * i)
        fill = colors[i % len(colors)]
        sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, top, cw, Inches(0.55))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        sh.line.fill.background(); sh.shadow.inherit = False
        tf = sh.text_frame; tf.word_wrap = False
        tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = label; p.font.size = Pt(SZ_L); p.font.color.rgb = WHITE
        p.font.bold = True; p.font.name = FN; p.alignment = PP_ALIGN.CENTER
        if p.runs: _set_ea(p.runs[0])
        t = tb(slide, x + Inches(0.1), top + Inches(0.6), cw - Inches(0.2), Inches(0.25))
        stxt(t.text_frame, sub, sz=SZ_M, color=TEXT_CAPTION, al=PP_ALIGN.CENTER)

def icon_label_row(slide, left, top, width, items):
    n = len(items); gap = int(width / n)
    for i, (icon, label, desc) in enumerate(items):
        cx = int(left + gap * i + gap / 2)
        sh = oval(slide, cx - Inches(0.35), top, Inches(0.7), Inches(0.7), SJTU_RED)
        stxt(sh.text_frame, icon, sz=SZ_H, color=WHITE, bold=True, al=PP_ALIGN.CENTER)
        t = tb(slide, cx - Inches(0.75), top + Inches(0.85), Inches(1.5), Inches(0.3))
        stxt(t.text_frame, label, sz=SZ_B, color=TEXT_PRIMARY, bold=True, al=PP_ALIGN.CENTER)
        t = tb(slide, cx - Inches(0.75), top + Inches(1.15), Inches(1.5), Inches(0.5))
        t.text_frame.word_wrap = True
        stxt(t.text_frame, desc, sz=SZ_L, color=TEXT_CAPTION, al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
#  PAGE 1: COVER
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(LAY_COVER)
for ph in s.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = ''
        p = ph.text_frame.paragraphs[0]
        r = p.add_run(); r.text = 'Claude Code 源码核心机制深度解析'
        r.font.size = Pt(40); r.font.color.rgb = WHITE; r.font.name = FN; r.font.bold = True; _set_ea(r)
        p2 = ph.text_frame.add_paragraph(); p2.space_before = Pt(20)
        r2 = p2.add_run(); r2.text = '从 System Prompt 到 MCP 协议，拆解 AI 编程 Agent 的工程细节'
        r2.font.size = Pt(20); r2.font.color.rgb = WHITE; r2.font.name = FN; _set_ea(r2)
    elif idx == 11:
        ph.text = ''
        p = ph.text_frame.paragraphs[0]
        r = p.add_run(); r.text = '2026年4月'
        r.font.size = Pt(16); r.font.color.rgb = WHITE; r.font.name = FN; _set_ea(r)


# ═══════════════════════════════════════════════════════════
#  PAGE 2: TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '报告框架', 2)

toc_groups = [
    ('核心引擎', SJTU_RED, [
        ('一', 'System Prompt', '03-04'),
        ('二', '工具系统', '05-10'),
        ('三', '目录树感知', '11'),
        ('四', 'Plan 模式', '12'),
    ]),
    ('智能管理', GOLD, [
        ('五', 'Context 压缩', '13-15'),
        ('六', 'Sub-Agent', '16-17'),
        ('七', '失败处理', '18'),
    ]),
    ('平台能力', RGBColor(0x8B, 0x5C, 0x3E), [
        ('八', 'Hooks 系统', '19'),
        ('九', 'CLAUDE.md 记忆', '20'),
        ('十', '权限与治理', '21-22'),
    ]),
    ('工程基础', TEXT_PRIMARY, [
        ('⑪', '状态持久化', '23-24'),
        ('⑫', 'MCP 协议', '25'),
        ('⑬', '预算管理', '26'),
        ('⑭', '总结对比', '27'),
    ]),
]

for gi, (group_title, accent, items) in enumerate(toc_groups):
    col = gi % 2; row = gi // 2
    gx = MARGIN + Inches(col * 5.95)
    gy = Inches(1.2) + Inches(row * 3.0)
    gh = Inches(2.7)
    card(s, gx, gy, Inches(5.65), gh, WHITE, WARM_GRAY_2)
    rect(s, gx, gy, Inches(5.65), Pt(3), accent)
    pill(s, gx + Inches(0.25), gy + Inches(0.2), Inches(1.8), Inches(0.32), group_title, bg=accent, fg=WHITE)
    for ii, (num, title, pg) in enumerate(items):
        iy = gy + Inches(0.72) + Inches(ii * 0.48)
        badge(s, gx + Inches(0.3), iy, num, size=Inches(0.36), sz=SZ_L, bg=accent)
        t = tb(s, gx + Inches(0.8), iy + Inches(0.02), Inches(3.0), Inches(0.3))
        stxt(t.text_frame, title, sz=SZ_B, color=TEXT_PRIMARY, bold=True)
        t = tb(s, gx + Inches(4.2), iy + Inches(0.02), Inches(1.0), Inches(0.3))
        stxt(t.text_frame, f'P.{pg}', sz=SZ_L, color=accent, al=PP_ALIGN.RIGHT, bold=True)


# ═══════════════════════════════════════════════════════════
#  PAGE 3: SYSTEM PROMPT INTRO
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '1. System Prompt：动态组装，而非静态模板', 3)

# Left card — diff
card(s, MARGIN, Inches(1.2), Inches(5.6), Inches(2.2), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(1.15), Inches(1.42), 'A', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(1.65), Inches(1.38), Inches(4), Inches(0.35))
stxt(t.text_frame, '核心差异', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(1.2))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '大多数框架：写死的静态文本，启动时原样注入',
    'Claude Code：运行时动态组装，由 buildEffectiveSystemPrompt 现场构建',
], sz=SZ_B, color=TEXT_BODY)

# Right card — 行为契约表
card(s, Inches(6.7), Inches(1.2), Inches(5.6), Inches(2.55), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(6.95), Inches(1.42), 'B', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(7.45), Inches(1.38), Inches(4.5), Inches(0.35))
stxt(t.text_frame, '默认 Prompt 的行为契约', sz=SZ_H, color=TEXT_PRIMARY, bold=True)

hdrs = ['规则', '说明']
data = [
    ['工具优先', '优先用 Read/Grep/Glob，而非 bash'],
    ['输出风格', '简洁直接，禁 emoji、填充词'],
    ['Memory', 'CLAUDE.md 自动发现路径'],
    ['Git 安全', 'force push / reset 需显式授权'],
]
table = make_tbl(s, 5, 2, Inches(7.0), Inches(1.88), Inches(5.0), Inches(1.7))
table.columns[0].width = Inches(1.5); table.columns[1].width = Inches(3.5)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 2)

# Bottom — 6 types injected
t = tb(s, MARGIN, Inches(4.0), CW, Inches(0.35))
stxt(t.text_frame, '运行时动态注入的 6 类内容', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
gold_div(s, MARGIN, Inches(4.35), CW)

items = [
    ('1', '工具描述', '遍历启用工具的 prompt()'),
    ('2', 'MCP 指令', '连接的 MCP 服务器说明'),
    ('3', 'Skill 索引', '已安装 skill 信息'),
    ('4', '环境信息', '平台/日期/工作目录'),
    ('5', 'ToolSearch', '延迟加载发现方式'),
    ('6', '用户配置', '6 层优先级合并'),
]
for i, (num, label, desc) in enumerate(items):
    x = MARGIN + Inches(i * 1.93)
    y = Inches(4.55)
    card(s, x, y, Inches(1.78), Inches(1.4), WHITE, WARM_GRAY_2)
    badge(s, x + Inches(0.65), y + Inches(0.1), num, size=Inches(0.38), sz=SZ_L)
    t = tb(s, x + Inches(0.1), y + Inches(0.55), Inches(1.58), Inches(0.3))
    stxt(t.text_frame, label, sz=SZ_B, color=TEXT_PRIMARY, bold=True, al=PP_ALIGN.CENTER)
    t = tb(s, x + Inches(0.05), y + Inches(0.85), Inches(1.68), Inches(0.4))
    t.text_frame.word_wrap = True
    stxt(t.text_frame, desc, sz=SZ_M, color=TEXT_CAPTION, al=PP_ALIGN.CENTER)

insight(s, MARGIN, Inches(6.2), CW,
        'Prompt 类型对比：Claude Code 动态组装（6层优先级） vs Codex/OpenCode/Gemini 静态模板')


# ═══════════════════════════════════════════════════════════
#  PAGE 4: SYSTEM PROMPT COMPARISON
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '1. System Prompt · 横向对比', 4)

hdrs = ['特性', 'Claude Code', 'Codex', 'OpenCode', 'Gemini-CLI']
data = [
    ['Prompt 类型', '动态组装（6层优先级）', '静态模板', '按模型选择静态文件', '静态模板'],
    ['工具描述注入', '每个工具自带 prompt()', '静态描述', '静态描述', '静态描述'],
    ['配置层级', '企业 > 用户 > 项目', '单层', '单层', '单层'],
    ['运行时感知', '平台/日期/目录/git', '目录', '无', '无'],
]
table = make_tbl(s, 5, 5, MARGIN, Inches(1.2), Inches(11.5), Inches(3.2))
table.columns[0].width = Inches(2.0)
for c in range(1, 5): table.columns[c].width = Inches(2.375)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 5)

chevron_flow(s, MARGIN, Inches(4.7), CW, [
    ('静态注入', '启动时固定'),
    ('按模型选择', '多文件切换'),
    ('动态组装', '运行时构建'),
    ('多层优先级', '企业级治理'),
])

insight(s, MARGIN, Inches(5.9), CW,
        '动态 Prompt 组装让 Claude Code 能根据运行上下文、启用工具、MCP 服务器实时调整行为')


# ═══════════════════════════════════════════════════════════
#  PAGE 5: TOOL SYSTEM OVERVIEW
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '2. 工具系统概览', 5)

stat_callout(s, MARGIN, Inches(1.2), Inches(2.8), Inches(1.15),
             '~45', '个工具', '工具总量', '分布在 40+ 子目录', num_sz=32)
stat_callout(s, Inches(4.0), Inches(1.2), Inches(2.8), Inches(1.15),
             '10', '并发上限', '最大并发数', '可通过环境变量覆盖', num_sz=32, accent=GOLD)
stat_callout(s, Inches(7.1), Inches(1.2), Inches(2.8), Inches(1.15),
             '4', '维度', '预算控制粒度', 'Token/成本/结果/轮次', num_sz=32)
stat_callout(s, Inches(10.2), Inches(1.2), Inches(2.15), Inches(1.15),
             '7', '种', 'Agent 模式', '含 Worktree 隔离', num_sz=32, accent=SJTU_DEEP)

icon_label_row(s, MARGIN, Inches(2.7), CW, [
    ('\u2713', '并发安全', 'isConcurrencySafe\n只读/写操作分批'),
    ('\u2630', '最大结果', 'maxResultSizeChars\n超限存磁盘'),
    ('\u26d4', '权限检查', 'checkPermissions\n三种决策路径'),
    ('\u23f3', '延迟加载', 'shouldDefer\n按需注入 Schema'),
])

quote_block(s, MARGIN, Inches(5.0), CW,
            '调度层在不了解工具内部实现的情况下，统一管理并发、权限和 token 预算',
            '设计目标')


# ═══════════════════════════════════════════════════════════
#  PAGE 6: TOOL CONCURRENCY
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '2. 工具并发调度：isConcurrencySafe', 6)

# Left card — rules
card(s, MARGIN, Inches(1.2), Inches(5.6), Inches(2.0), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(1.15), Inches(1.42), 'A', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(1.65), Inches(1.38), Inches(4), Inches(0.35))
stxt(t.text_frame, '并发规则', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(1.0))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '\u2713 只读操作（Glob/Grep/Read/WebSearch）\u2192 可并发',
    '\u2717 写操作（Edit/Write/Bash）\u2192 必须串行',
    '最大并发数默认 10（环境变量可覆盖）',
], sz=SZ_B, color=TEXT_BODY)

# Right card — batch example
card(s, Inches(6.7), Inches(1.2), Inches(5.6), Inches(2.0), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(6.95), Inches(1.42), 'B', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(7.45), Inches(1.38), Inches(4.5), Inches(0.35))
stxt(t.text_frame, '分批执行示例', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(7.0), Inches(1.9), Inches(5.0), Inches(1.0))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, '模型输出：[Glob, Grep, Read, FileEdit, Glob, Read]', sz=SZ_L, color=TEXT_CAPTION)
apara(tf, '\u2460 [Glob, Grep, Read] \u2192 Promise.all \u5e76\u53d1', sz=SZ_B, color=TEXT_BODY, sb=Pt(10))
apara(tf, '\u2461 [FileEdit] \u2192 \u4e32\u884c\u6267\u884c', sz=SZ_B, color=TEXT_BODY)
apara(tf, '\u2462 [Glob, Read] \u2192 Promise.all \u5e76\u53d1', sz=SZ_B, color=TEXT_BODY)

# Step flow
step_flow(s, Inches(2.0), Inches(3.5), CW - Inches(2.0), [
    ('\u2460', '只读并发'), ('\u2461', '写操作串行'), ('\u2462', '只读并发'),
])

insight(s, MARGIN, Inches(4.65), CW,
        '模型需在单次回复中发出多个只读调用，才能利用并发 — 隐式要求模型具备批次意识')

# Bottom comparison
hdrs = ['特性', 'Claude Code', 'Codex', 'OpenCode', 'Gemini-CLI']
data = [
    ['工具并发', '自动分批', '\u2717', '手动 batch', '自动'],
    ['延迟加载', '\u2713 shouldDefer', '\u2717', '\u2717', '\u2717'],
    ['结果大小', '超限存磁盘', '截断（首尾）', '无', '截断'],
]
table = make_tbl(s, 4, 5, MARGIN, Inches(5.1), Inches(11.5), Inches(1.5))
table.columns[0].width = Inches(2.0)
for c in range(1, 5): table.columns[c].width = Inches(2.375)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 4, 5)


# ═══════════════════════════════════════════════════════════
#  PAGE 7: TOOL DEFERRED LOADING
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '2. 工具延迟加载：shouldDefer + ToolSearch', 7)

# Problem-solution
card(s, MARGIN, Inches(1.2), Inches(5.6), Inches(1.5), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(1.2), Pt(3), Inches(1.5), WARN_AMBER)
t = tb(s, Inches(1.2), Inches(1.28), Inches(4.8), Inches(0.3))
stxt(t.text_frame, '\u26a0 \u95ee\u9898', sz=SZ_H, color=WARN_AMBER, bold=True)
t = tb(s, Inches(1.2), Inches(1.65), Inches(4.8), Inches(0.7))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, '工具越多 \u2192 Schema 越占 context \u2192 Token 浪费严重（尤其 MCP 工具）', sz=SZ_B, color=TEXT_BODY)

card(s, Inches(6.7), Inches(1.2), Inches(5.6), Inches(1.5), WHITE, border=WARM_GRAY_2)
rect(s, Inches(6.7), Inches(1.2), Pt(3), Inches(1.5), SUCCESS)
t = tb(s, Inches(7.0), Inches(1.28), Inches(4.8), Inches(0.3))
stxt(t.text_frame, '\u2713 \u89e3\u6cd5', sz=SZ_H, color=SUCCESS, bold=True)
t = tb(s, Inches(7.0), Inches(1.65), Inches(5.0), Inches(0.7))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, 'shouldDefer: true \u2192 初始只发空壳（无参数描述）\n模型调用 ToolSearch 搜索 \u2192 注入完整 Schema', sz=SZ_B, color=TEXT_BODY)

# Search scoring
t = tb(s, MARGIN, Inches(2.95), Inches(4), Inches(0.35))
stxt(t.text_frame, 'ToolSearch 搜索评分权重', sz=SZ_H, color=TEXT_PRIMARY, bold=True)

hdrs = ['匹配来源', '权重']
data = [
    ['searchHint（3~10 词）', '4 分'],
    ['工具名', '2 分'],
    ['完整 prompt 描述', '1 分'],
]
table = make_tbl(s, 4, 2, MARGIN, Inches(3.35), Inches(4.5), Inches(1.7))
table.columns[0].width = Inches(3.0); table.columns[1].width = Inches(1.5)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 4, 2)

# Prefix cache
card(s, Inches(5.8), Inches(2.95), Inches(6.5), Inches(2.1), WHITE, WARM_GRAY_2)
rect(s, Inches(5.8), Inches(2.95), Pt(3), Inches(2.1), GOLD)
t = tb(s, Inches(6.1), Inches(3.05), Inches(5.8), Inches(0.3))
stxt(t.text_frame, '避免破坏 Prefix Cache', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(6.1), Inches(3.5), Inches(5.8), Inches(1.3))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '早期方案：动态插入已发现工具 \u2192 cache 每次失效 \u2717',
    '当前方案：通过独立 attachment 发送 \u2192 消息 prefix 不变 \u2713',
], sz=SZ_B, color=TEXT_BODY)

insight(s, MARGIN, Inches(5.3), CW,
        'Prefix Cache 命中率直接影响 API 成本和延迟 — 这是延迟加载最精妙的工程权衡')


# ═══════════════════════════════════════════════════════════
#  PAGE 8: TOOL RESULT SIZE & PERMISSIONS
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '2. 工具结果大小控制 & 权限检查', 8)

card(s, MARGIN, Inches(1.2), Inches(5.6), Inches(3.0), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(1.15), Inches(1.42), '1', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(1.65), Inches(1.38), Inches(4), Inches(0.35))
stxt(t.text_frame, '结果大小控制', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(2.0))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '每个工具声明 maxResultSizeChars',
    '超出 \u2192 自动持久化到磁盘，模型收到路径引用',
    '例外：Read 工具设为 Infinity（否则"读\u2192路径\u2192再读"死循环）',
], sz=SZ_B, color=TEXT_BODY)

card(s, Inches(6.7), Inches(1.2), Inches(5.6), Inches(3.0), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(6.95), Inches(1.42), '2', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(7.45), Inches(1.38), Inches(4.5), Inches(0.35))
stxt(t.text_frame, '权限检查', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(7.0), Inches(1.85), Inches(5.0), Inches(0.3))
stxt(t.text_frame, '每个工具独立 checkPermissions() \u2192 三种结果：', sz=SZ_B, color=TEXT_BODY)

perms = [
    ('\u2713', '自动放行', 'alwaysAllow 规则匹配', SUCCESS),
    ('?', '询问用户', '无规则匹配', GOLD),
    ('\u26d4', '直接拦截', 'alwaysDeny / Hooks 拦截', SJTU_RED),
]
for i, (icon, label, desc, color) in enumerate(perms):
    y = Inches(2.3) + Inches(i * 0.6)
    sh = oval(s, Inches(7.2), y, Inches(0.4), Inches(0.4), color)
    stxt(sh.text_frame, icon, sz=SZ_B, color=WHITE, bold=True, al=PP_ALIGN.CENTER)
    t = tb(s, Inches(7.75), y + Inches(0.02), Inches(1.5), Inches(0.35))
    stxt(t.text_frame, label, sz=SZ_B, color=TEXT_PRIMARY, bold=True)
    t = tb(s, Inches(9.3), y + Inches(0.02), Inches(3.0), Inches(0.35))
    stxt(t.text_frame, desc, sz=SZ_B, color=TEXT_CAPTION)

insight(s, MARGIN, Inches(4.5), CW,
        '工具结果超限存磁盘 + Read 设 Infinity 例外 — 细节中体现了对"模型行为"的深入理解')


# ═══════════════════════════════════════════════════════════
#  PAGE 9: TOOL CATALOG
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '2. 工具分类总览', 9)

hdrs = ['类别', '工具', '说明']
data = [
    ['文件操作', 'Read / Edit / Write / Glob / Grep', '常用 shell 命令抽象为独立工具'],
    ['Shell', 'Bash', '持久化 shell 会话，跨调用保留环境'],
    ['Multi-Agent', 'Agent / SendMessage / TeamCreate', '子 Agent 入口 + Swarms 协作'],
    ['规划', 'EnterPlanMode / ExitPlanMode', '权限层面只读约束'],
    ['任务追踪', 'TaskCreate / TaskUpdate / TodoWrite', '任务与进度管理'],
    ['搜索网络', 'WebSearch / WebFetch / ToolSearch', '网络搜索 + 延迟工具发现'],
    ['MCP', 'MCPTool / ListMcpResources', '外部工具/资源接入'],
    ['高级', 'LSP / Worktree / Cron / REPL', '代码导航 / 沙箱 / 定时'],
]
table = make_tbl(s, 9, 3, MARGIN, Inches(1.2), Inches(11.5), Inches(5.3))
table.columns[0].width = Inches(1.8); table.columns[1].width = Inches(4.5); table.columns[2].width = Inches(5.2)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 9, 3)


# ═══════════════════════════════════════════════════════════
#  PAGE 10: TOOL COMPARISON
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '2. 工具系统横向对比', 10)

hdrs = ['特性', 'Claude Code', 'Codex', 'OpenCode', 'Gemini-CLI']
data = [
    ['工具并发', '自动分批', '\u2717', '手动 batch', '自动'],
    ['延迟加载', '\u2713 shouldDefer', '\u2717', '\u2717', '\u2717'],
    ['结果大小限制', '超限存磁盘', '截断（首尾保留）', '无', '截断'],
    ['LSP 工具', '\u2713', '\u2717', '\u2713', '\u2717'],
    ['语义代码搜索', '\u2717', '\u2717', '\u2713（Exa Code）', '\u2717'],
]
table = make_tbl(s, 6, 5, MARGIN, Inches(1.2), Inches(11.5), Inches(3.5))
table.columns[0].width = Inches(2.0)
for c in range(1, 5): table.columns[c].width = Inches(2.375)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 6, 5)

insight(s, MARGIN, Inches(5.0), CW,
        'Claude Code 在工具调度精细度上远超同类框架，延迟加载和磁盘持久化是独有设计')


# ═══════════════════════════════════════════════════════════
#  PAGE 11: REPO TREE AWARENESS
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '3. 仓库目录树感知', 11)

hdrs = ['框架', '方式', '说明']
data = [
    ['Claude Code', '不注入目录树，注入 git 状态', '每轮更新，按需探索'],
    ['Codex', '自动生成 2 层目录树', '注入 user prompt'],
    ['OpenCode', '硬编码禁用', '&& false 强制跳过'],
    ['Gemini-CLI', '不注入目录树', '注入 git 工作流指引'],
]
table = make_tbl(s, 5, 3, MARGIN, Inches(1.2), Inches(11.5), Inches(2.5))
table.columns[0].width = Inches(2.2); table.columns[1].width = Inches(4.3); table.columns[2].width = Inches(5.0)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 3)

t = tb(s, MARGIN, Inches(3.95), Inches(5), Inches(0.35))
stxt(t.text_frame, 'Claude Code 的 git 上下文（每轮注入）', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
gold_div(s, MARGIN, Inches(4.3), CW)

items = [
    ('\u2295', '分支名', '当前分支 + HEAD'),
    ('\u2261', 'Commit 记录', '最近几条提交'),
    ('\u270e', '工作区变更', 'git status（最多 2000 字符）'),
]
for i, (icon, label, desc) in enumerate(items):
    x = MARGIN + Inches(i * 3.95)
    card(s, x, Inches(4.5), Inches(3.65), Inches(1.2), WHITE, WARM_GRAY_2)
    sh = oval(s, x + Inches(0.2), Inches(4.65), Inches(0.4), Inches(0.4), SJTU_RED)
    stxt(sh.text_frame, icon, sz=SZ_B, color=WHITE, bold=True, al=PP_ALIGN.CENTER)
    t = tb(s, x + Inches(0.75), Inches(4.6), Inches(2.5), Inches(0.3))
    stxt(t.text_frame, label, sz=SZ_B, color=TEXT_PRIMARY, bold=True)
    t = tb(s, x + Inches(0.75), Inches(4.95), Inches(2.5), Inches(0.5))
    t.text_frame.word_wrap = True
    stxt(t.text_frame, desc, sz=SZ_L, color=TEXT_CAPTION)

quote_block(s, MARGIN, Inches(5.9), CW,
            '"当前改了哪些文件" 比 "目录里有哪些文件"更有决策价值 — 目录结构由模型主动探索，不占固定 token 预算',
            '设计判断')


# ═══════════════════════════════════════════════════════════
#  PAGE 12: PLAN MODE
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '4. Plan 模式', 12)

card(s, MARGIN, Inches(1.2), CW, Inches(0.5), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(1.2), Pt(3), Inches(0.5), GOLD)
t = tb(s, MARGIN + Pt(16), Inches(1.2), CW - Pt(20), Inches(0.5))
stxt(t.text_frame, '核心区别：Claude Code 用权限系统层面约束（mode = plan），写操作直接拦截，不靠模型自觉', sz=SZ_B, color=TEXT_BODY)
try:
    bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 'ctr')
except: pass

# Three triggers
step_flow(s, Inches(1.5), Inches(1.95), CW - Inches(1.2), [
    ('\u2460', '模型调用 EnterPlanMode'), ('\u2461', '启动参数 --mode plan'), ('\u2462', '用户 UI 手动切换'),
])

# Exit flow
t = tb(s, MARGIN, Inches(2.8), CW, Inches(0.35))
stxt(t.text_frame, '退出与审批流程', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
gold_div(s, MARGIN, Inches(3.15), CW)

chevron_flow(s, MARGIN, Inches(3.35), CW, [
    ('规划执行', '只读工具'),
    ('ExitPlanMode', '写入 .claude/plans/'),
    ('UI 弹出审批', '用户必须批准'),
    ('执行实施', '写操作解锁'),
])

card(s, MARGIN, Inches(4.3), Inches(5.6), Inches(1.6), WHITE, WARM_GRAY_2)
rect(s, MARGIN, Inches(4.3), Pt(3), Inches(1.6), SJTU_RED)
t = tb(s, Inches(1.2), Inches(4.4), Inches(4.8), Inches(0.3))
stxt(t.text_frame, '唯一强制用户介入的环节', sz=SZ_B, color=SJTU_RED, bold=True)
t = tb(s, Inches(1.2), Inches(4.8), Inches(4.8), Inches(0.8))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, 'UI 弹出审批对话框 \u2192 用户必须手动批准规划方案后才能进入执行阶段', sz=SZ_B, color=TEXT_BODY)

card(s, Inches(6.7), Inches(4.3), Inches(5.6), Inches(1.6), GOLD_PALE, border=WARN_AMBER, shadow=False)
rect(s, Inches(6.7), Inches(4.3), Pt(3), Inches(1.6), WARN_AMBER)
t = tb(s, Inches(7.0), Inches(4.4), Inches(4.8), Inches(0.3))
stxt(t.text_frame, '\u26a0 限制', sz=SZ_B, color=WARN_AMBER, bold=True)
t = tb(s, Inches(7.0), Inches(4.8), Inches(5.0), Inches(0.8))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, '禁止在子 Agent 中使用 — 子 Agent 无法弹出 UI 审批，会永远无法退出 Plan 模式', sz=SZ_B, color=TEXT_BODY)

insight(s, MARGIN, Inches(6.15), CW,
        'Plan 模式本质：用系统权限做 guardrail，而非靠 prompt 指令约束模型行为')


# ═══════════════════════════════════════════════════════════
#  PAGE 13: CONTEXT COMPRESSION — FIVE LAYERS
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '5. Context 压缩管理 — 五层递进机制', 13)

card(s, MARGIN, Inches(1.2), CW, Inches(0.5), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(1.2), Pt(3), Inches(0.5), GOLD)
t = tb(s, MARGIN + Pt(16), Inches(1.2), CW - Pt(20), Inches(0.5))
stxt(t.text_frame, '阈值与当前模型 context window 动态绑定 | 公式：context_window - 输出保留(20K) - buffer(13K)', sz=SZ_B, color=TEXT_BODY)
try:
    bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 'ctr')
except: pass

layers = [
    ('1', '工具结果预算', '超限存磁盘（每轮执行）', '轻量'),
    ('2', '历史片段截断', '规则打分，删低分消息', '规则'),
    ('3', '微压缩', '利用 cache_edits 服务端清理', '保 Cache'),
    ('4', '上下文折叠', '旧内容 \u2192 摘要，近期保留原始', 'LLM'),
    ('5', '完整摘要压缩', 'fork 子 Agent 调 LLM 生成摘要', '重量级'),
]
for i, (num, label, desc, tag) in enumerate(layers):
    y = Inches(1.95) + Inches(i * 0.9)
    card(s, MARGIN, y, CW, Inches(0.7), WHITE, WARM_GRAY_2)
    badge(s, Inches(1.1), y + Inches(0.12), num, size=Inches(0.42), sz=SZ_L)
    t = tb(s, Inches(1.65), y + Inches(0.08), Inches(2.5), Inches(0.3))
    stxt(t.text_frame, label, sz=SZ_H, color=TEXT_PRIMARY, bold=True)
    t = tb(s, Inches(4.3), y + Inches(0.08), Inches(6.0), Inches(0.55))
    t.text_frame.word_wrap = True
    stxt(t.text_frame, desc, sz=SZ_B, color=TEXT_BODY)
    sh = rrect(s, Inches(10.5), y + Inches(0.15), Inches(1.5), Inches(0.33), GOLD_PALE, border=GOLD, bw=Pt(0.5))
    stxt(sh.text_frame, tag, sz=SZ_L, color=SJTU_DEEP, bold=True, al=PP_ALIGN.CENTER)

insight(s, MARGIN, Inches(6.65), CW,
        '从轻量到昂贵逐层兜底 — 200K window 下 ~167K 自动压缩 / 160K 警告 / 177K 硬性拦截')


# ═══════════════════════════════════════════════════════════
#  PAGE 14: MICRO COMPACT
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '5. Context 压缩 · 第 3 层 — 微压缩 (microCompact)', 14)

card(s, MARGIN, Inches(1.2), CW, Inches(0.5), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(1.2), Pt(3), Inches(0.5), GOLD)
t = tb(s, MARGIN + Pt(16), Inches(1.2), CW - Pt(20), Inches(0.5))
stxt(t.text_frame, '核心思路：不修改本地消息 \u2192 利用 API cache_edits 服务端清理旧工具结果', sz=SZ_B, color=TEXT_BODY)
try:
    bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 'ctr')
except: pass

hdrs = ['模式', '条件', '操作']
data = [
    ['时间触发', '距上次消息 > 60 分钟（cache 已过期）', '直接修改本地消息，替换为 [cleared]'],
    ['热缓存模式', 'cache 仍有效', '通过 cache_edits 服务端注意力屏蔽'],
]
table = make_tbl(s, 3, 3, MARGIN, Inches(1.95), Inches(11.5), Inches(1.5))
table.columns[0].width = Inches(2.0); table.columns[1].width = Inches(4.5); table.columns[2].width = Inches(5.0)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 3, 3)

quote_block(s, MARGIN, Inches(3.7), CW,
            'cache_edits 不是删除 token，而是将 attention mask 置 0 — 序列位置编码不变 \u2192 KV 缓存全部有效 \u2192 Cache 持续命中',
            '反直觉的关键')

t = tb(s, MARGIN, Inches(5.45), Inches(3), Inches(0.3))
stxt(t.text_frame, '可清理的工具', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
gold_div(s, MARGIN, Inches(5.75), CW)

tools = ['Read', 'Bash', 'Grep', 'Glob', 'WebSearch', 'WebFetch', 'FileEdit', 'FileWrite']
for i, tool in enumerate(tools):
    x = MARGIN + Inches(i * 1.44)
    pill(s, x, Inches(5.95), Inches(1.3), Inches(0.28), tool)


# ═══════════════════════════════════════════════════════════
#  PAGE 15: CONTEXT LAYERS 4 & 5 + COMPARISON
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '5. Context 压缩 · 第 4 & 5 层 + 横向对比', 15)

card(s, MARGIN, Inches(1.2), Inches(5.6), Inches(2.2), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(1.15), Inches(1.42), '4', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(1.65), Inches(1.38), Inches(4), Inches(0.35))
stxt(t.text_frame, '上下文折叠 (contextCollapse)', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(1.2))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '历史分组 \u2192 旧分组归档为摘要，近期保留原始粒度',
    '触发：context 用量 ~90% 准备，95% 阻塞触发',
    '与 autoCompact 互斥（避免竞争覆盖）',
], sz=SZ_B, color=TEXT_BODY)

card(s, Inches(6.7), Inches(1.2), Inches(5.6), Inches(2.2), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(6.95), Inches(1.42), '5', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(7.45), Inches(1.38), Inches(4.5), Inches(0.35))
stxt(t.text_frame, '完整摘要压缩 (autoCompact)', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(7.0), Inches(1.9), Inches(5.0), Inches(1.2))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    'Fork 子 Agent 调 LLM 生成完整对话摘要',
    '压缩后：[compact_boundary] + [摘要] + [尾部消息]',
    '重新注入：CLAUDE.md / MCP 指令 / Skill 列表',
], sz=SZ_B, color=TEXT_BODY)

hdrs = ['框架', '触发方式', '是否调 LLM', '特点']
data = [
    ['Claude Code', '响应式，每轮检查', '\u2713 fork 子 Agent', '五层递进 + 附件恢复'],
    ['Codex', '预防性 + 响应式', '\u2713', '保留最近用户消息'],
    ['Gemini-CLI', '主动式', '\u2713', '保留最新 30% 历史'],
    ['OpenCode', '响应式', '\u2713', '全量替换'],
]
table = make_tbl(s, 5, 4, MARGIN, Inches(3.65), Inches(11.5), Inches(2.6))
table.columns[0].width = Inches(2.0); table.columns[1].width = Inches(2.8)
table.columns[2].width = Inches(2.5); table.columns[3].width = Inches(4.2)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 4)


# ═══════════════════════════════════════════════════════════
#  PAGE 16: SUB-AGENT SYSTEM
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '6. Sub-Agent 系统 — 7 种执行模式', 16)

card(s, MARGIN, Inches(1.2), CW, Inches(0.5), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(1.2), Pt(3), Inches(0.5), GOLD)
t = tb(s, MARGIN + Pt(16), Inches(1.2), CW - Pt(20), Inches(0.5))
stxt(t.text_frame, '统一入口 AgentTool \u2192 通过参数组合路由到 7 种执行模式', sz=SZ_B, color=TEXT_BODY)
try:
    bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 'ctr')
except: pass

hdrs = ['模式', '触发条件', '特点']
data = [
    ['同步前台', '默认', '阻塞等待结果'],
    ['异步后台', 'run_in_background: true', '立即返回 ID，轮询结果'],
    ['自动转后台', '运行 > 120 秒', '自动切换，避免阻塞'],
    ['Worktree 隔离', "isolation: 'worktree'", '独立 git 副本'],
    ['远端执行', "isolation: 'remote'", '云端运行（内部功能）'],
    ['Fork 模式', '实验性', '继承父 Agent 完整历史'],
    ['Teammate', 'Swarms 模式', '独立 tmux，双向通信'],
]
table = make_tbl(s, 8, 3, MARGIN, Inches(1.95), Inches(11.5), Inches(4.6))
table.columns[0].width = Inches(2.5); table.columns[1].width = Inches(4.0); table.columns[2].width = Inches(5.0)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 8, 3)


# ═══════════════════════════════════════════════════════════
#  PAGE 17: SUB-AGENT TYPES & CONTEXT
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '6. Sub-Agent · 内置类型 & Context 共享', 17)

hdrs = ['类型', '工具集', '适用场景']
data = [
    ['general-purpose', '所有工具（除 AgentTool）', '通用复杂任务'],
    ['Explore', '只读工具', '代码库探索'],
    ['Plan', '只读 + ExitPlanMode', '规划阶段'],
    ['claude-code-guide', 'Read / Grep / WebSearch', 'Claude Code 使用答疑'],
]
table = make_tbl(s, 5, 3, MARGIN, Inches(1.2), Inches(11.5), Inches(2.5))
table.columns[0].width = Inches(2.8); table.columns[1].width = Inches(4.2); table.columns[2].width = Inches(4.5)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 3)

card(s, MARGIN, Inches(4.0), Inches(5.6), Inches(2.2), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(1.15), Inches(4.22), 'A', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(1.65), Inches(4.18), Inches(4), Inches(0.35))
stxt(t.text_frame, '普通模式共享', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(1.2), Inches(4.7), Inches(4.8), Inches(1.2))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '克隆文件缓存',
    '工具结果预算',
    '权限上下文',
    'MCP 连接',
], sz=SZ_B, color=TEXT_BODY)

card(s, Inches(6.7), Inches(4.0), Inches(5.6), Inches(2.2), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(6.95), Inches(4.22), 'B', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(7.45), Inches(4.18), Inches(4.5), Inches(0.35))
stxt(t.text_frame, 'Fork 模式额外继承', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(7.0), Inches(4.7), Inches(5.0), Inches(1.2))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '完整对话历史',
    '字节级相同的 System Prompt',
    '\u2192 保证 Prefix Cache 命中',
], sz=SZ_B, color=TEXT_BODY)

insight(s, MARGIN, Inches(6.45), CW,
        '支持用户自定义 Agent（YAML 定义文件，自动发现注册）')


# ═══════════════════════════════════════════════════════════
#  PAGE 18: FAILURE HANDLING
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '7. 失败处理机制', 18)

# Tool errors
card(s, MARGIN, Inches(1.2), Inches(5.6), Inches(1.55), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(1.15), Inches(1.42), '1', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(1.65), Inches(1.38), Inches(4), Inches(0.35))
stxt(t.text_frame, '工具执行错误', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(1.2), Inches(1.85), Inches(4.8), Inches(0.7))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '错误以 is_error: true 返回给模型',
    '批次内单个失败不中断其他工具',
    '无失败计数预算，不会主动终止',
], sz=SZ_L, color=TEXT_BODY)

# API errors
card(s, Inches(6.7), Inches(1.2), Inches(5.6), Inches(1.55), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(6.95), Inches(1.42), '2', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(7.45), Inches(1.38), Inches(4.5), Inches(0.35))
stxt(t.text_frame, 'API 错误恢复', sz=SZ_H, color=TEXT_PRIMARY, bold=True)

hdrs = ['错误类型', '处理方式']
data = [
    ['输出 token 超限', '自动重试 \u2264 3 次'],
    ['请求过长', '先 autoCompact \u2192 再逐条删'],
    ['网络失败', '指数退避重试'],
]
table = make_tbl(s, 4, 2, Inches(7.0), Inches(1.88), Inches(5.0), Inches(0.75))
table.columns[0].width = Inches(2.0); table.columns[1].width = Inches(3.0)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 4, 2)

# Permission escalation
card(s, MARGIN, Inches(3.1), CW, Inches(1.65), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(3.1), Pt(3), Inches(1.65), SJTU_RED)
t = tb(s, Inches(1.2), Inches(3.18), Inches(10), Inches(0.3))
stxt(t.text_frame, '权限拒绝的渐进式升级（Claude Code 特有）', sz=SZ_H, color=SJTU_DEEP, bold=True)
t = tb(s, Inches(1.2), Inches(3.55), Inches(10), Inches(0.9))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, '连续拒绝 \u2265 3 次 或 累计拒绝 \u2265 20 次 \u2192 从"自动拒绝"切换到"询问用户" \u2192 安全阀，避免死循环', sz=SZ_B, color=TEXT_BODY)

# No loop detection warning
card(s, MARGIN, Inches(5.0), CW, Inches(1.3), WHITE, border=WARN_AMBER, shadow=False)
rect(s, MARGIN, Inches(5.0), Pt(3), Inches(1.3), WARN_AMBER)
t = tb(s, Inches(1.2), Inches(5.08), Inches(10), Inches(0.3))
stxt(t.text_frame, '\u26a0 无内置死循环检测', sz=SZ_H, color=WARN_AMBER, bold=True)
t = tb(s, Inches(1.2), Inches(5.45), Inches(10), Inches(0.7))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    'OpenCode：相同工具+参数连续 3 次 \u2192 询问用户',
    'Gemini-CLI：注入恢复 prompt \u2192 60 秒强制终止',
    'Claude Code：仅靠 autoCompact + 用户 ESC 兜底',
], sz=SZ_L, color=TEXT_BODY)


# ═══════════════════════════════════════════════════════════
#  PAGE 19: HOOKS SYSTEM
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '8. Hooks 系统（Claude Code 独有）', 19)

card(s, MARGIN, Inches(1.2), CW, Inches(0.5), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(1.2), Pt(3), Inches(0.5), GOLD)
t = tb(s, MARGIN + Pt(16), Inches(1.2), CW - Pt(20), Inches(0.5))
stxt(t.text_frame, '定位：将工具调用流程从"黑盒"变为可扩展平台 | 24 种 Hook 事件覆盖 6 大类别', sz=SZ_B, color=TEXT_BODY)
try:
    bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 'ctr')
except: pass

hdrs = ['能力', 'Hook 返回字段']
data = [
    ['\u26d4 拦截工具调用', "decision: 'block'"],
    ['\u270e 修改工具输入', 'updatedInput'],
    ['\U0001f4ce 注入上下文', 'additionalContext'],
    ['\U0001f4dd 修改工具输出', 'updatedMCPToolOutput'],
    ['\U0001f4ac 替换初始消息', 'initialUserMessage'],
    ['\u23f9 终止会话', 'continue: false'],
    ['\U0001f511 自动化权限', 'allow / deny'],
]
table = make_tbl(s, 8, 2, MARGIN, Inches(1.95), Inches(5.5), Inches(3.8))
table.columns[0].width = Inches(2.5); table.columns[1].width = Inches(3.0)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 8, 2)

# Right side: categories + config
card(s, Inches(6.7), Inches(1.95), Inches(5.6), Inches(2.0), WHITE, WARM_GRAY_2)
rect(s, Inches(6.7), Inches(1.95), Pt(3), Inches(2.0), SJTU_RED)
t = tb(s, Inches(7.0), Inches(2.05), Inches(5.0), Inches(0.3))
stxt(t.text_frame, '事件类别覆盖', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
cats = ['生命周期', '工具', '权限', 'Sub-Agent', '用户交互', '压缩', '任务', '系统', 'MCP']
for i, cat in enumerate(cats):
    row = i // 3; col = i % 3
    x = Inches(7.0) + Inches(col * 1.7)
    y = Inches(2.5) + Inches(row * 0.42)
    pill(s, x, y, Inches(1.5), Inches(0.28), cat)

card(s, Inches(6.7), Inches(4.2), Inches(5.6), Inches(1.55), WHITE, WARM_GRAY_2)
rect(s, Inches(6.7), Inches(4.2), Pt(3), Inches(1.55), GOLD)
t = tb(s, Inches(7.0), Inches(4.3), Inches(5.0), Inches(0.3))
stxt(t.text_frame, '三层配置优先级', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(7.0), Inches(4.7), Inches(5.0), Inches(0.8))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, '企业管理策略 > 用户级 > 项目级', sz=SZ_B, color=TEXT_BODY, bold=True)
apara(tf, '--no-hooks 一键禁用', sz=SZ_L, color=TEXT_CAPTION, sb=Pt(8))

insight(s, MARGIN, Inches(6.0), CW,
        'Hooks 系统使 Claude Code 从"工具"升级为"平台" — Codex / OpenCode / Gemini 均无此能力')


# ═══════════════════════════════════════════════════════════
#  PAGE 20: CLAUDE.MD MEMORY
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '9. CLAUDE.md 记忆系统', 20)

# Directory discovery
card(s, MARGIN, Inches(1.2), Inches(5.6), Inches(3.2), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(1.15), Inches(1.42), 'A', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(1.65), Inches(1.38), Inches(4), Inches(0.35))
stxt(t.text_frame, '多层目录递归发现', sz=SZ_H, color=TEXT_PRIMARY, bold=True)

paths = [
    ('~/.claude/CLAUDE.md', '全局用户级（最低优先级）'),
    ('<project_root>/CLAUDE.md', '项目级'),
    ('<current_dir>/CLAUDE.md', '目录级'),
    ('<parent_dirs>/CLAUDE.md', '递归向上（最高优先级）'),
]
for i, (path, desc) in enumerate(paths):
    y = Inches(1.95) + Inches(i * 0.55)
    sh = rrect(s, Inches(1.2), y, Inches(4.0), Inches(0.4), WARM_GRAY_1, border=WARM_GRAY_2, bw=Pt(0.5))
    stxt(sh.text_frame, path, sz=SZ_L, color=TEXT_PRIMARY, bold=True, al=PP_ALIGN.CENTER)
    t = tb(s, Inches(1.2), y + Inches(0.4), Inches(4.0), Inches(0.2))
    stxt(t.text_frame, desc, sz=SZ_M, color=TEXT_CAPTION, al=PP_ALIGN.CENTER)

# Typical uses
card(s, Inches(6.7), Inches(1.2), Inches(5.6), Inches(3.2), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(6.95), Inches(1.42), 'B', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(7.45), Inches(1.38), Inches(4.5), Inches(0.35))
stxt(t.text_frame, '典型用途', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(7.0), Inches(1.9), Inches(5.0), Inches(2.2))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '项目约束：代码风格、禁止修改的目录、命名规则',
    '常用命令：build / test / lint',
    '架构说明：模块职责、依赖关系',
    '跨会话记忆：模型可主动写入，实现跨会话学习',
], sz=SZ_B, color=TEXT_BODY)

# Comparison table
hdrs = ['框架', '记忆机制', '作用范围']
data = [
    ['Claude Code', 'CLAUDE.md（多层递归）', '全局 / 项目 / 目录'],
    ['Codex', 'AGENTS.md', '项目级'],
    ['Gemini-CLI', 'GEMINI.md', '项目级'],
    ['OpenCode', '无', '\u2014'],
]
table = make_tbl(s, 5, 3, MARGIN, Inches(4.7), Inches(11.5), Inches(2.0))
table.columns[0].width = Inches(2.2); table.columns[1].width = Inches(4.3); table.columns[2].width = Inches(5.0)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 3)


# ═══════════════════════════════════════════════════════════
#  PAGE 21: PERMISSIONS
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '10. 权限与治理系统', 21)

hdrs = ['模式', '行为', '场景']
data = [
    ['default', '每次弹确认框', '交互式 IDE'],
    ['auto', 'AI 分类器自动判断', 'CI / 自动化'],
    ['plan', '只读放行，写操作拦截', '规划阶段'],
    ['bypassPermissions', '跳过所有检查', '完全自动化'],
]
table = make_tbl(s, 5, 3, MARGIN, Inches(1.2), Inches(11.5), Inches(2.5))
table.columns[0].width = Inches(2.8); table.columns[1].width = Inches(4.2); table.columns[2].width = Inches(4.5)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 3)

gold_div(s, MARGIN, Inches(3.9), CW)

t = tb(s, MARGIN, Inches(4.05), Inches(4), Inches(0.35))
stxt(t.text_frame, '静态规则：Allow / Deny / Ask', sz=SZ_H, color=TEXT_PRIMARY, bold=True)

rules = [
    ('alwaysAllow', 'Bash(git *)', '直接放行', SUCCESS),
    ('alwaysDeny', 'Bash(rm -rf *)', '始终拒绝', SJTU_RED),
    ('alwaysAsk', 'Bash(deploy *)', '即使 auto 也询问', WARN_AMBER),
]
for i, (rule, example, action, color) in enumerate(rules):
    y = Inches(4.5) + Inches(i * 0.65)
    card(s, MARGIN, y, CW, Inches(0.5), WHITE, WARM_GRAY_2, shadow=False)
    rect(s, MARGIN, y, Pt(3), Inches(0.5), color)
    t = tb(s, Inches(1.2), y + Inches(0.05), Inches(2.0), Inches(0.35))
    stxt(t.text_frame, rule, sz=SZ_B, color=TEXT_PRIMARY, bold=True)
    t = tb(s, Inches(3.3), y + Inches(0.05), Inches(4.0), Inches(0.35))
    stxt(t.text_frame, example, sz=SZ_B, color=TEXT_CAPTION)
    t = tb(s, Inches(8.0), y + Inches(0.05), Inches(3.5), Inches(0.35))
    stxt(t.text_frame, '\u2192 ' + action, sz=SZ_B, color=color, bold=True)

insight(s, MARGIN, Inches(6.65), CW,
        '三层优先级：企业管理策略 > 用户级 > 项目级')


# ═══════════════════════════════════════════════════════════
#  PAGE 22: AI CLASSIFIER
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '10. AI 安全分类器（Auto Mode）', 22)

card(s, MARGIN, Inches(1.2), CW, Inches(0.5), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(1.2), Pt(3), Inches(0.5), GOLD)
t = tb(s, MARGIN + Pt(16), Inches(1.2), CW - Pt(20), Inches(0.5))
stxt(t.text_frame, '不是规则匹配 \u2192 是一次真正的 AI 模型调用（Claude Opus） | 输入：工具调用描述 + 对话历史紧凑编码', sz=SZ_B, color=TEXT_BODY)
try:
    bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 'ctr')
except: pass

# Two-stage
step_flow(s, Inches(2.5), Inches(1.9), CW - Inches(3.0), [
    ('\u2460', 'Stage 1 快速判断'), ('\u2461', 'Stage 2 深度推理'),
])

hdrs = ['阶段', 'max_tokens', '行为']
data = [
    ['Stage 1（快速）', '64', '输出 yes/no \u2192 "不阻止"立即返回'],
    ['Stage 2（思考）', '4096', '<thinking> 链式推理 \u2192 降低误报率'],
]
table = make_tbl(s, 3, 3, MARGIN, Inches(2.7), Inches(11.5), Inches(1.3))
table.columns[0].width = Inches(2.5); table.columns[1].width = Inches(2.0); table.columns[2].width = Inches(7.0)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 3, 3)

t = tb(s, MARGIN, Inches(4.2), Inches(3), Inches(0.35))
stxt(t.text_frame, '优化措施', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
gold_div(s, MARGIN, Inches(4.55), CW)

t = tb(s, Inches(1.2), Inches(4.75), Inches(10.5), Inches(1.2))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '35+ 只读工具走白名单，直接放行，不调分类器',
    '两阶段都利用 prompt cache 复用前缀',
    '分类器不可用时默认 Fail-closed（阻止，而非降级放行）',
    'Worktree 沙箱：子 Agent 在独立 git 副本中操作，主工作区不受影响',
], sz=SZ_B, color=TEXT_BODY)

insight(s, MARGIN, Inches(6.65), CW,
        '两阶段分类 + Fail-closed 默认策略 — 安全性设计远超同类框架')


# ═══════════════════════════════════════════════════════════
#  PAGE 23: STATE PERSISTENCE
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '11. 状态持久化与会话恢复', 23)

card(s, MARGIN, Inches(1.2), Inches(5.6), Inches(2.5), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(1.15), Inches(1.42), 'A', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(1.65), Inches(1.38), Inches(4), Inches(0.35))
stxt(t.text_frame, 'Session 存储', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(1.5))
tf = t.text_frame; tf.word_wrap = True
bullets(tf, [
    '完整 transcript \u2192 JSONL 格式',
    '路径：~/.claude/projects/<hash>/',
    '记录类型：消息 / compact_boundary / 摘要 / 文件替换 / worktree',
    'System prompt 不存储（恢复时重新组装）',
], sz=SZ_B, color=TEXT_BODY)

card(s, Inches(6.7), Inches(1.2), Inches(5.6), Inches(2.5), WHITE, WARM_GRAY_2, gold_top=True)
badge(s, Inches(6.95), Inches(1.42), 'B', size=Inches(0.38), sz=SZ_L)
t = tb(s, Inches(7.45), Inches(1.38), Inches(4.5), Inches(0.35))
stxt(t.text_frame, 'Compact 后存储结构', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(7.0), Inches(1.9), Inches(5.0), Inches(0.6))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, '[旧历史消息...] [compact_boundary] [摘要] [新消息...]', sz=SZ_L, color=TEXT_PRIMARY, bold=True)
apara(tf, 'logicalParentUuid \u2192 指向压缩前最后一条消息', sz=SZ_L, color=TEXT_CAPTION, sb=Pt(8))
apara(tf, '文件 > 50MB \u2192 跳过 boundary 之前，只读后半部分', sz=SZ_L, color=TEXT_CAPTION, sb=Pt(6))

# Sub-agent sidechain
card(s, Inches(7.0), Inches(3.1), Inches(5.0), Inches(0.5), WARM_GRAY_1, WARM_GRAY_2, shadow=False)
rect(s, Inches(7.0), Inches(3.1), Pt(3), Inches(0.5), GOLD)
t = tb(s, Inches(7.3), Inches(3.1), Inches(4.5), Inches(0.5))
stxt(t.text_frame, '子 Agent 独立存储：subagents/agent-<id>.jsonl', sz=SZ_L, color=TEXT_BODY)
try:
    bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 'ctr')
except: pass

# Resume flow
t = tb(s, MARGIN, Inches(4.0), Inches(6), Inches(0.35))
stxt(t.text_frame, '跨会话恢复流程（/resume）— 四步重建', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
gold_div(s, MARGIN, Inches(4.35), CW)

steps = [
    ('1', '确定读取范围', '找最后一个 compact_boundary \u2192 只加载其后的消息'),
    ('2', '重建对话链', '按 parentUuid 构建 DAG \u2192 过滤孤立分支 \u2192 线性序列'),
    ('3', '恢复应用状态', 'content-replacement / 文件历史 / TodoWrite / worktree'),
    ('4', '重注入动态内容', 'System prompt / CLAUDE.md / MCP / Skill 重新注入'),
]
for i, (num, label, desc) in enumerate(steps):
    x = MARGIN + Inches(i * 2.92)
    y = Inches(4.55)
    card(s, x, y, Inches(2.7), Inches(1.8), WHITE, WARM_GRAY_2)
    badge(s, x + Inches(1.0), y + Inches(0.12), num, size=Inches(0.5))
    t = tb(s, x + Inches(0.1), y + Inches(0.7), Inches(2.5), Inches(0.3))
    stxt(t.text_frame, label, sz=SZ_B, color=TEXT_PRIMARY, bold=True, al=PP_ALIGN.CENTER)
    t = tb(s, x + Inches(0.1), y + Inches(1.05), Inches(2.5), Inches(0.65))
    t.text_frame.word_wrap = True
    stxt(t.text_frame, desc, sz=SZ_M, color=TEXT_CAPTION, al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
#  PAGE 24: MCP PROTOCOL
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '12. MCP 协议集成', 24)

card(s, MARGIN, Inches(1.2), CW, Inches(0.5), GOLD_PALE, border=GOLD, shadow=False)
rect(s, MARGIN, Inches(1.2), Pt(3), Inches(0.5), GOLD)
t = tb(s, MARGIN + Pt(16), Inches(1.2), CW - Pt(20), Inches(0.5))
stxt(t.text_frame, 'Claude Code 是四个框架中唯一完整原生支持 MCP 的', sz=SZ_B, color=TEXT_BODY, bold=True)
try:
    bodyPr = t.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None: bodyPr.set('anchor', 'ctr')
except: pass

icon_label_row(s, MARGIN, Inches(2.0), CW, [
    ('\u2295', '动态工具扩展', '第三方注册 mcp__<server>__<tool>\n共享全部调度机制'),
    ('\u2261', '资源访问', '文件 / 数据库 / API\n结构化数据读取'),
    ('\U0001f511', '认证 & 交互', 'OAuth 认证流程 +\nElicitation 用户输入'),
])

# Config example card
card(s, MARGIN, Inches(4.3), CW, Inches(2.1), WARM_GRAY_1, border=WARM_GRAY_2)
rect(s, MARGIN, Inches(4.3), Pt(3), Inches(2.1), GOLD)
t = tb(s, Inches(1.2), Inches(4.38), Inches(3), Inches(0.3))
stxt(t.text_frame, '配置示例（settings.json）', sz=SZ_H, color=TEXT_PRIMARY, bold=True)
t = tb(s, Inches(1.2), Inches(4.75), Inches(10.5), Inches(1.4))
tf = t.text_frame; tf.word_wrap = True
stxt(tf, 'mcpServers: { "filesystem": { command: "npx", args: ["@modelcontextprotocol/server-filesystem", "/tmp"] }, '
     '"github": { command: "npx", args: ["@modelcontextprotocol/server-github"], env: { GITHUB_TOKEN: "..." } } }',
     sz=SZ_L, color=TEXT_BODY)

insight(s, MARGIN, Inches(6.6), CW,
        '扩展边界由 MCP 生态而非 Claude Code 本身决定 — 这是平台化设计的核心体现')


# ═══════════════════════════════════════════════════════════
#  PAGE 25: BUDGET MANAGEMENT
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '13. 预算管理 — 四维度独立控制', 25)

hdrs = ['维度', '机制', '说明']
data = [
    ['Token 预算', 'output_config.task_budget', '整个 agentic turn 的 token 总量上限'],
    ['成本预算', 'maxBudgetUsd', '单次会话最大美元成本'],
    ['工具结果预算', 'maxResultSizeChars', '单次工具结果字符上限，超限存磁盘'],
    ['轮次预算', 'maxTurns', 'Agent 最大迭代次数，防失控'],
]
table = make_tbl(s, 5, 3, MARGIN, Inches(1.2), Inches(11.5), Inches(2.5))
table.columns[0].width = Inches(2.0); table.columns[1].width = Inches(3.8); table.columns[2].width = Inches(5.7)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 3)

gold_div(s, MARGIN, Inches(3.9), CW)

hdrs = ['框架', 'Token', '成本', '工具结果', '轮次']
data = [
    ['Claude Code', '\u2713', '\u2713', '\u2713 存磁盘', '\u2713'],
    ['Codex', '\u2717', '\u2717', '截断', '\u2713'],
    ['OpenCode', '\u2717', '\u2717', '\u2717', '\u2717'],
    ['Gemini-CLI', '\u2717', '\u2717', '截断', '\u2713'],
]
table = make_tbl(s, 5, 5, MARGIN, Inches(4.15), Inches(11.5), Inches(2.1))
table.columns[0].width = Inches(2.0)
for c in range(1, 5): table.columns[c].width = Inches(2.375)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 5, 5)

insight(s, MARGIN, Inches(6.5), CW,
        'Claude Code 是唯一实现四维度预算独立控制的框架 — 精细化资源管理降低生产事故风险')


# ═══════════════════════════════════════════════════════════
#  PAGE 26: DESIGN PHILOSOPHY SUMMARY
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '总结：Claude Code 的核心设计哲学', 26)

philosophies = [
    ('\u2699', '动态而非静态', [
        'System Prompt 运行时组装',
        '工具 Schema 延迟加载',
        '压缩阈值与 context window 动态绑定',
    ]),
    ('\u26d4', '约束落在系统层', [
        'Plan 模式 \u2192 权限系统层面拦截写操作',
        'AI 分类器 \u2192 两阶段分类 + 渐进升级',
        'Worktree \u2192 文件系统级隔离沙箱',
    ]),
    ('\u2295', '平台化而非工具化', [
        'Hooks \u2192 24 种事件，全生命周期可扩展',
        'MCP \u2192 动态接入任意第三方能力',
        '自定义 Agent \u2192 YAML 定义即注册',
    ]),
    ('\u2630', '精细化资源管理', [
        '五层 Context 压缩递进机制',
        '四维度预算独立控制',
        '完整的状态持久化与无损会话恢复',
    ]),
]
for i, (icon, title, items) in enumerate(philosophies):
    col = i % 2; row = i // 2
    x = MARGIN + Inches(col * 5.85)
    y = Inches(1.2) + Inches(row * 2.8)
    card(s, x, y, Inches(5.55), Inches(2.5), WHITE, WARM_GRAY_2, gold_top=True)
    sh = oval(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5), SJTU_RED)
    stxt(sh.text_frame, icon, sz=SZ_H, color=WHITE, bold=True, al=PP_ALIGN.CENTER)
    t = tb(s, x + Inches(0.85), y + Inches(0.22), Inches(4), Inches(0.4))
    stxt(t.text_frame, title, sz=SZ_H, color=TEXT_PRIMARY, bold=True)
    t = tb(s, x + Inches(0.25), y + Inches(0.8), Inches(5.0), Inches(1.4))
    tf = t.text_frame; tf.word_wrap = True
    bullets(tf, items, sz=SZ_B, color=TEXT_BODY)


# ═══════════════════════════════════════════════════════════
#  PAGE 27: COMPARISON OVERVIEW
# ═══════════════════════════════════════════════════════════
s = add_content_slide()
sjtu_header(s, '横向对比总览', 27)

hdrs = ['能力维度', 'Claude Code', 'Codex', 'OpenCode', 'Gemini-CLI']
data = [
    ['System Prompt', '动态组装', '静态', '按模型选择', '静态'],
    ['工具数量', '~45', '~5', '~10', '~10'],
    ['工具并发', '自动分批', '\u2717', '手动 batch', '自动'],
    ['延迟加载', '\u2713', '\u2717', '\u2717', '\u2717'],
    ['Context 压缩', '五层递进', '单层', '单层', '单层'],
    ['Sub-Agent', '7 种模式', '3 角色', '2 层', '1 层'],
    ['Hooks 系统', '24 种事件', '\u2717', '\u2717', '\u2717'],
    ['MCP 支持', '\u2713 完整', '\u2717', '\u2717', '\u2717'],
    ['预算维度', '4 维', '1 维', '0 维', '1 维'],
    ['权限系统', 'AI 分类器', '静态规则', '静态规则', '策略引擎'],
]
table = make_tbl(s, 11, 5, MARGIN, Inches(1.2), Inches(11.5), Inches(5.2))
table.columns[0].width = Inches(2.0)
for c in range(1, 5): table.columns[c].width = Inches(2.375)
tbl_hdr(table, hdrs); tbl_data(table, data); style_tbl(table, 11, 5)


# ═══════════════════════════════════════════════════════════
#  PAGE 28: THANK YOU
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(LAY_END)
for ph in s.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 11:
        ph.text = ''
        r = ph.text_frame.paragraphs[0].add_run()
        r.text = 'Thank You'
        r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
        r.font.name = FN; _set_ea(r)
        p2 = ph.text_frame.add_paragraph()
        r2 = p2.add_run()
        r2.text = 'Claude Code 源码核心机制深度解析'
        r2.font.size = Pt(18); r2.font.color.rgb = GOLD_LIGHT
        r2.font.name = FN; _set_ea(r2)
        p3 = ph.text_frame.add_paragraph()
        r3 = p3.add_run()
        r3.text = '对比框架：Claude Code / Codex / OpenCode / Gemini-CLI  |  2026年4月'
        r3.font.size = Pt(14); r3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        r3.font.name = FN; _set_ea(r3)


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
output = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Claude_Code_源码解析_SJTU.pptx')
prs.save(output)
print(f'\u2728 SJTU Premium PPT \u5df2\u751f\u6210: {output}')
print(f'   \u5171 {len(prs.slides)} \u9875 | 16:9 \u5bbd\u5c4f | SJTU \u6a21\u677f\u98ce\u683c')
print(f'   \u2714 Wine Red + Gold \u914d\u8272 \u00b7 Claude Code \u6e90\u7801\u89e3\u6790')
