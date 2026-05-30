import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .constants import (
    GOLD,
    GOLD_PALE,
    SJTU_DEEP,
    SJTU_RED,
    WHITE,
)
from .schema import Column, Deck, Layout, Slide


SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.82
CONTENT_TOP_IN = 1.12
CONTENT_BOTTOM_IN = 6.63
CONTENT_W_IN = SLIDE_W_IN - MARGIN_IN * 2
FONT = "Microsoft YaHei"
MATH_FONT = "Cambria Math"
TEXT_PRIMARY = "1A1A2E"
TEXT_BODY = "3D3D4E"
TEXT_CAPTION = "6B7280"
WARM_GRAY_1 = "F5F4F2"
WARM_GRAY_2 = "EDECEA"
WARM_GRAY_3 = "D2D0CD"
BROWN = "8B5C3E"
OFF_WHITE = "FAFAFB"
DATE_RE = re.compile(r"(?:20\d{2}[./-]\d{1,2}[./-]\d{1,2}|20\d{2}\s*年|三年前|上周)")
TABLE_RULE_RE = re.compile(r"^\|?\s*:?-{3,}:?\s*(?:\|\s*:?-{3,}:?\s*)+\|?$")
KEY_VALUE_RE = re.compile(r"^([^：:]{1,14})[：:]\s*(.+)$")
FORMULA_RE = re.compile(r"(\$[^$]+\$|\\\(|\\\[|[=∑∫√≤≥≈θλμσᵢŷ])")
LATEX_REPLACEMENTS = {
    r"\theta": "θ",
    r"\lambda": "λ",
    r"\mu": "μ",
    r"\sigma": "σ",
    r"\alpha": "α",
    r"\beta": "β",
    r"\hat{y}": "ŷ",
    r"\sum": "∑",
    r"\sqrt": "√",
    r"\le": "≤",
    r"\ge": "≥",
    r"\approx": "≈",
    r"\times": "×",
    r"\cdot": "·",
}
TEMPLATE_CANDIDATES = (
    "AIPPT_TEMPLATE_PPTX",
    "AIPPT_TEMPLATE_PPTX_PATH",
)
TEMPLATE_LAYOUTS = {
    Layout.COVER: ("封面-01", 0),
    Layout.TOC: ("1_目录-1", 3),
    Layout.THANKS: ("封底01", 12),
}
TEMPLATE_CONTENT_LAYOUT = "常规样式（1）"
COVER_TITLE_LIMIT = 18
COVER_SUBTITLE_LIMIT = 28


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def build_pptx(deck: Deck, output_path: str | Path) -> Path:
    preserve_template_cover = bool(deck.slides and deck.slides[0].layout == Layout.COVER)
    prs, use_template, template_cover_ready = _new_presentation(preserve_template_cover=preserve_template_cover)

    for idx, slide in enumerate(deck.slides, start=1):
        if use_template and template_cover_ready and idx == 1 and slide.layout == Layout.COVER:
            _render_template_cover(prs.slides[0], slide)
            continue
        render_slide(prs, slide, idx, use_template=use_template)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def _new_presentation(preserve_template_cover: bool = False) -> tuple[Presentation, bool, bool]:
    template_path = _template_path()
    if template_path is not None:
        prs = Presentation(template_path)
        template_cover_ready = preserve_template_cover and len(prs.slides) > 0
        _remove_existing_slides(prs, keep_first=template_cover_ready)
        return prs, True, template_cover_ready
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs, False, False


def _template_path() -> Path | None:
    candidates = [Path(value) for key in TEMPLATE_CANDIDATES if (value := os.getenv(key))]
    if repo_root := os.getenv("AIPPT_REPO_ROOT"):
        candidates.append(Path(repo_root) / "docs" / "SJTU PPT 模板" / "SJTU 模板.pptx")
    candidates.extend(
        parent / "docs" / "SJTU PPT 模板" / "SJTU 模板.pptx"
        for parent in Path(__file__).resolve().parents
    )
    candidates.append(Path.cwd() / "assets" / "SJTU 模板.pptx")
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    return None


def _remove_existing_slides(prs: Presentation, keep_first: bool = False) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    if keep_first:
        slide_ids = slide_ids[1:]
    for slide_id in slide_ids:
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def render_slide(
    prs: Presentation,
    slide_data: Slide,
    page_num: int,
    use_template: bool = False,
) -> None:
    if use_template:
        render_template_slide(prs, slide_data, page_num)
        return

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if slide_data.layout == Layout.COVER:
        render_cover(slide, slide_data)
    elif slide_data.layout == Layout.THANKS:
        render_thanks(slide, slide_data)
    elif slide_data.layout == Layout.TOC:
        render_header(slide, slide_data.title)
        render_toc(slide, slide_data)
        render_footer(slide, page_num)
    else:
        render_header(slide, slide_data.title)
        render_body(slide, slide_data)
        render_footer(slide, page_num)


def render_template_slide(prs: Presentation, slide_data: Slide, page_num: int) -> None:
    layout = _template_layout(prs, slide_data.layout)
    slide = prs.slides.add_slide(layout)
    if slide_data.layout == Layout.COVER:
        _render_template_cover(slide, slide_data)
    elif slide_data.layout == Layout.TOC:
        _render_template_toc(slide, slide_data)
    elif slide_data.layout == Layout.THANKS:
        _render_template_thanks(slide, slide_data)
    else:
        _render_template_content(slide, slide_data)


def _template_layout(prs: Presentation, layout: Layout):
    if layout in TEMPLATE_LAYOUTS:
        name, fallback_index = TEMPLATE_LAYOUTS[layout]
    else:
        name, fallback_index = TEMPLATE_CONTENT_LAYOUT, 6
    for candidate in prs.slide_layouts:
        if candidate.name == name:
            return candidate
    return prs.slide_layouts[min(fallback_index, len(prs.slide_layouts) - 1)]


def _render_template_cover(slide, slide_data: Slide) -> None:
    title, subtitle = _cover_title_and_subtitle(slide_data)
    _set_placeholder_lines(
        slide,
        0,
        [title],
        size=41 if len(title) <= 12 else 35,
        color=WHITE,
        bold=True,
    )
    _set_placeholder_lines(
        slide,
        11,
        [subtitle] if subtitle else ["AI PPT 生成工作台"],
        size=17,
        color=WHITE,
        bold=True,
    )


def _render_template_content(slide, slide_data: Slide) -> None:
    _set_placeholder_lines(slide, 11, [slide_data.title], size=25, color=WHITE, bold=True)
    _remove_placeholder(slide, 14)
    items = _clean_items(slide_data.bullets or [c.heading for c in slide_data.columns if c.heading])
    if not items:
        items = ["请补充这一页的核心结论。"]
    render_body(
        slide,
        Slide(
            title=slide_data.title,
            layout=slide_data.layout,
            visual=slide_data.visual,
            proof=slide_data.proof,
            support=slide_data.support,
            bullets=items,
            columns=slide_data.columns,
            items=slide_data.items,
            table=slide_data.table,
            insight=slide_data.insight,
        ),
    )


def _render_template_toc(slide, slide_data: Slide) -> None:
    items = _clean_items(slide_data.bullets)[:8]
    x0, y0 = 4.35, 1.55
    for idx, item in enumerate(items, start=1):
        y = y0 + (idx - 1) * 0.55
        box = slide.shapes.add_textbox(Inches(x0), Inches(y), Inches(7.1), Inches(0.34))
        p = box.text_frame.paragraphs[0]
        p.text = f"{idx:02d}  {item}"
        _style_paragraph(p, 17, TEXT_BODY, bold=True)


def _render_template_thanks(slide, slide_data: Slide) -> None:
    placeholder = _placeholder(slide, 11)
    if placeholder is not None:
        placeholder.left = Inches(3.15)
        placeholder.top = Inches(2.95)
        placeholder.width = Inches(2.54)
        placeholder.height = Inches(0.86)
        placeholder.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_placeholder_lines(
        slide,
        11,
        [slide_data.title or "谢谢"],
        size=36,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _cover_title_and_subtitle(slide_data: Slide) -> tuple[str, str]:
    raw_title = (slide_data.title or "AIPPT").strip()
    title, title_tail = _split_cover_title(raw_title)
    subtitle_candidates = [title_tail, *slide_data.subtitle.splitlines()]
    subtitle = ""
    for candidate in subtitle_candidates:
        compact = _compact_cover_subtitle(candidate, title)
        if compact:
            subtitle = compact
            break
    return _trim(title, COVER_TITLE_LIMIT), _trim(subtitle, COVER_SUBTITLE_LIMIT)


def _split_cover_title(title: str) -> tuple[str, str]:
    for delimiter in ("：", ":", "——", "—", " - "):
        if delimiter not in title:
            continue
        head, tail = title.split(delimiter, 1)
        head = head.strip()
        tail = tail.strip()
        if 3 <= len(head) <= COVER_TITLE_LIMIT + 4 and tail:
            return head, tail
    return title, ""


def _compact_cover_subtitle(text: str, title: str) -> str:
    text = re.sub(r"^\s*副标题\s*[：:]\s*", "", text.strip())
    text = re.sub(r"\s+", " ", text)
    if not text or text == title:
        return ""
    if any(skip in text for skip in ("可继续编辑大纲", "一句话需求自动规划")):
        return ""
    match = re.match(r"用\s*([0-9一二三四五六七八九十]+)\s*页讲清楚", text)
    if match:
        return f"{match.group(1)} 页速览"
    if len(text) > COVER_SUBTITLE_LIMIT and "——" in text:
        text = text.split("——", 1)[0].strip()
    return text


def _set_placeholder_lines(
    slide,
    idx: int,
    lines: list[str],
    size: float,
    color: str,
    bold: bool = False,
    font_name: str = FONT,
    align: PP_ALIGN | None = None,
) -> None:
    placeholder = _placeholder(slide, idx)
    if placeholder is None:
        return
    frame = placeholder.text_frame
    frame.clear()
    frame.word_wrap = True
    for line_idx, line in enumerate(lines):
        p = frame.paragraphs[0] if line_idx == 0 else frame.add_paragraph()
        p.text = line
        _style_paragraph(p, size, color, bold=bold, font_name=font_name, align=align)
        p.line_spacing = 1.06


def _placeholder(slide, idx: int):
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx == idx:
            return placeholder
    return None


def _clear_placeholder(slide, idx: int) -> None:
    placeholder = _placeholder(slide, idx)
    if placeholder is not None and placeholder.has_text_frame:
        placeholder.text_frame.clear()


def _remove_placeholder(slide, idx: int) -> None:
    placeholder = _placeholder(slide, idx)
    if placeholder is None:
        return
    placeholder._element.getparent().remove(placeholder._element)


def render_cover(slide, slide_data: Slide) -> None:
    _fill_background(slide, SJTU_RED)
    _rect(slide, 9.65, 0, 3.68, 7.5, SJTU_DEEP)
    _rect(slide, 0, 6.92, 13.333, 0.08, GOLD)
    _rect(slide, 0.86, 1.12, 0.09, 4.95, GOLD)
    _decorative_rule(slide, 1.05, 1.18, 4.55)

    title = slide_data.title or "AIPPT"
    title_box = slide.shapes.add_textbox(Inches(1.18), Inches(1.72), Inches(7.95), Inches(1.45))
    frame = title_box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = title
    _style_paragraph(p, 34 if len(title) <= 22 else 29, WHITE, bold=True)
    p.line_spacing = 1.05

    subtitle_lines = [line.strip() for line in slide_data.subtitle.splitlines() if line.strip()]
    subtitle_lines = subtitle_lines[:4] or ["AI PPT 生成工作台"]
    meta = slide.shapes.add_textbox(Inches(1.2), Inches(3.45), Inches(7.8), Inches(1.55))
    meta_frame = meta.text_frame
    meta_frame.word_wrap = True
    for idx, line in enumerate(subtitle_lines):
        p = meta_frame.paragraphs[0] if idx == 0 else meta_frame.add_paragraph()
        p.text = _trim(line, 54)
        _style_paragraph(p, 15 if idx == 0 else 12.5, WHITE, bold=(idx == 0))
        p.space_after = Pt(5)

    badge = _round_rect(slide, 1.2, 5.55, 2.2, 0.36, GOLD_PALE, border=GOLD, radius=0.14)
    _shape_text(badge, "SJTU Wine Red + Gold", 9, SJTU_DEEP, bold=True, align=PP_ALIGN.CENTER)


def render_thanks(slide, slide_data: Slide) -> None:
    _fill_background(slide, SJTU_RED)
    _rect(slide, 0, 6.92, 13.333, 0.08, GOLD)
    _rect(slide, 9.8, 0, 3.53, 7.5, SJTU_DEEP)
    box = slide.shapes.add_textbox(Inches(1.12), Inches(2.7), Inches(7.7), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = slide_data.title or "谢谢"
    _style_paragraph(p, 42, WHITE, bold=True)
    _decorative_rule(slide, 1.16, 3.88, 3.35)
    note = slide.shapes.add_textbox(Inches(1.15), Inches(4.12), Inches(6.2), Inches(0.4))
    p = note.text_frame.paragraphs[0]
    p.text = "AI PPT 生成工作台"
    _style_paragraph(p, 13, WHITE)


def render_header(slide, title: str) -> None:
    _fill_background(slide, OFF_WHITE)
    _rect(slide, 0, 0, 13.333, 0.78, SJTU_RED)
    _rect(slide, 0, 0.78, 13.333, 0.035, GOLD)
    _rect(slide, 11.35, 0, 1.98, 0.78, SJTU_DEEP)
    box = slide.shapes.add_textbox(Inches(0.78), Inches(0.18), Inches(10.4), Inches(0.42))
    p = box.text_frame.paragraphs[0]
    p.text = title
    _style_paragraph(p, _title_size(title), WHITE, bold=True)


def render_toc(slide, slide_data: Slide) -> None:
    bullets = _clean_items(slide_data.bullets)
    _section_label(slide, "目录", "从问题到方法，再到可执行的工作流")
    for idx, item in enumerate(bullets[:8], start=1):
        col = 0 if idx <= 4 else 1
        row = (idx - 1) % 4
        left = MARGIN_IN + col * 5.95
        top = 1.45 + row * 1.18
        _number_badge(slide, left, top + 0.07, f"{idx:02d}", color=SJTU_RED)
        title = slide.shapes.add_textbox(Inches(left + 0.58), Inches(top), Inches(4.75), Inches(0.38))
        p = title.text_frame.paragraphs[0]
        p.text = item
        _style_paragraph(p, 17, TEXT_PRIMARY, bold=True)
        _rect(slide, left + 0.58, top + 0.48, 4.55, 0.012, WARM_GRAY_3)


def render_body(slide, slide_data: Slide) -> None:
    items = _clean_items(slide_data.bullets or [c.heading for c in slide_data.columns if c.heading])
    if not items:
        items = ["请补充这一页的核心结论。"]

    visual = slide_data.visual or ""
    support = slide_data.support or slide_data.proof
    if visual == "metric_strip":
        _render_metric_strip(slide, items, support)
    elif visual == "milestone_timeline":
        _render_milestone_timeline(slide, items)
    elif visual == "project_showcase":
        _render_project_showcase(slide, items)
    elif visual == "media_explain":
        _render_media_explain(slide, items, support)
    elif visual == "stat_callout":
        _render_stat_callouts(slide, items, support)
    elif visual == "quote_block":
        _render_quote_block_slide(slide, items, support)
    elif visual == "concept_diagram":
        _render_concept_diagram(slide, items, support)
    elif visual == "example_walkthrough":
        _render_example_walkthrough(slide, items)
    elif slide_data.layout == Layout.TABLE or visual == "table":
        _render_table_slide(slide, slide_data.table, items)
    elif slide_data.layout in {Layout.TWO_COLUMN, Layout.COMPARISON} or visual == "two_column":
        _render_column_cards(slide, slide_data.columns, items, count=2)
    elif slide_data.layout == Layout.THREE_COLUMN or visual == "three_column":
        _render_column_cards(slide, slide_data.columns, items, count=3)
    elif slide_data.layout == Layout.HORIZONTAL and visual != "process":
        _render_horizontal_items(slide, slide_data.items, items)
    elif slide_data.layout == Layout.SUMMARY or visual == "summary":
        _render_summary(slide, items)
    elif slide_data.layout == Layout.SECTION:
        _render_section_break(slide, slide_data.title, items)
    elif visual == "timeline" or _looks_like_timeline(items):
        _render_timeline(slide, items)
    elif visual == "process" or _looks_like_process(slide_data.title, items):
        _render_process(slide, items)
    elif visual == "fact_grid":
        _render_fact_grid(slide, items)
    elif visual == "rich_cards" or _looks_like_rich_cards(items):
        _render_rich_card_grid(slide, items)
    elif _looks_like_fact_grid(items):
        _render_fact_grid(slide, items)
    else:
        _render_card_grid(slide, items)

    if slide_data.insight and not (visual == "stat_callout" and support):
        _insight(slide, slide_data.insight)


def _render_column_cards(slide, columns, items: list[str], count: int) -> None:
    lead, rest = _split_lead(items)
    if lead:
        _lead_callout(slide, lead, compact=True)
        top = 2.12
        height = 3.72
    else:
        rest = items
        top = 1.35
        height = 4.7

    if not columns:
        columns = _fallback_columns(rest, count)

    if count == 2:
        positions = [(0.92, top, 5.55, height), (6.82, top, 5.55, height)]
    else:
        positions = [
            (0.92, top, 3.64, height),
            (4.84, top, 3.64, height),
            (8.76, top, 3.64, height),
        ]

    for idx, column in enumerate(columns[:count], start=1):
        left, card_top, width, card_h = positions[idx - 1]
        accent = SJTU_RED if idx % 2 else GOLD
        _round_rect(slide, left, card_top, width, card_h, WHITE, border=WARM_GRAY_2, radius=0.035, shadow=True)
        _rect(slide, left, card_top, width, 0.05, accent)
        _number_badge(slide, left + 0.24, card_top + 0.24, str(idx), size=0.38, color=accent)

        heading = column.heading or f"要点 {idx}"
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.74),
            Inches(card_top + 0.2),
            Inches(width - 1.02),
            Inches(0.42),
        )
        title_box.text_frame.word_wrap = True
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_box.text_frame.paragraphs[0]
        p.text = _trim(heading, 24 if width >= 5 else 16)
        _style_paragraph(p, 14 if width >= 5 else 13, TEXT_PRIMARY, bold=True)

        points = column.bullets or [_trim(heading, 44)]
        _draw_card_points(
            slide,
            left + 0.32,
            card_top + 0.86,
            width - 0.62,
            card_h - 1.08,
            points,
            compact=count == 3,
        )


def _fallback_columns(items: list[str], count: int) -> list[Column]:
    columns: list[Column] = []
    for idx in range(count):
        chunk = items[idx::count] or [f"要点 {idx + 1}"]
        heading, body = _split_key_value(chunk[0])
        if body:
            points = _split_card_points(body)
        else:
            heading = heading if len(heading) <= 16 else f"要点 {idx + 1}"
            points = chunk[:4]
        columns.append(Column(heading=heading, bullets=points))
    return columns


def _render_horizontal_items(slide, horizontal_items, items: list[str]) -> None:
    lead, rest = _split_lead(items)
    if lead:
        _lead_callout(slide, lead, compact=True)
        top = 2.36
    else:
        rest = items
        top = 1.72

    if horizontal_items:
        process_items = [
            f"{item.heading}：{item.desc}" if item.desc else item.heading
            for item in horizontal_items[:5]
        ]
    else:
        process_items = rest[:5] or items[:5]

    left = 0.9
    width = 11.55
    gap = width / max(len(process_items), 1)
    colors = [SJTU_RED, GOLD, BROWN, TEXT_PRIMARY, SJTU_RED]
    for idx, item in enumerate(process_items, start=1):
        x = left + (idx - 1) * gap
        card_w = max(gap - 0.22, 2.0)
        _process_card(slide, x, top, card_w, 2.25, idx, item, colors[(idx - 1) % len(colors)])

    if len(items) > len(process_items) + 1:
        _insight(slide, items[-1], top=5.82)


def _render_summary(slide, items: list[str]) -> None:
    lead, rest = _split_lead(items)
    _lead_callout(slide, lead or "最后可以带走三句话。", compact=True)
    cards = (rest or items)[:3]
    positions = [
        (0.92, 2.35, 3.64, 3.05),
        (4.84, 2.35, 3.64, 3.05),
        (8.76, 2.35, 3.64, 3.05),
    ]
    for idx, item in enumerate(cards, start=1):
        left, top, width, height = positions[idx - 1]
        _rich_content_card(slide, left, top, width, height, idx, item)


def _render_section_break(slide, title: str, items: list[str]) -> None:
    _rect(slide, 0.9, 2.12, 0.08, 2.1, GOLD)
    box = slide.shapes.add_textbox(Inches(1.18), Inches(2.08), Inches(10.6), Inches(0.85))
    p = box.text_frame.paragraphs[0]
    p.text = title
    _style_paragraph(p, 31 if len(title) <= 18 else 26, TEXT_PRIMARY, bold=True)
    note = items[0] if items else ""
    if note:
        note_box = slide.shapes.add_textbox(Inches(1.22), Inches(3.12), Inches(9.4), Inches(0.7))
        note_box.text_frame.word_wrap = True
        p = note_box.text_frame.paragraphs[0]
        p.text = _trim(note, 96)
        _style_paragraph(p, 15, TEXT_BODY)


def _render_table_slide(slide, table_data, items: list[str]) -> None:
    if table_data and table_data.headers and table_data.rows:
        headers = table_data.headers[:4]
        rows = [row[: len(headers)] for row in table_data.rows[:5]]
    else:
        headers, rows = _fallback_table(items)

    rows_count = min(6, len(rows) + 1)
    cols_count = min(4, len(headers))
    table_height = min(4.45, 0.46 * rows_count + 0.2)
    shape = slide.shapes.add_table(
        rows_count,
        cols_count,
        Inches(0.92),
        Inches(1.42),
        Inches(11.55),
        Inches(table_height),
    )
    table = shape.table
    for col_idx in range(cols_count):
        table.columns[col_idx].width = Inches(11.55 / cols_count)

    for col_idx, header in enumerate(headers[:cols_count]):
        _style_table_cell(table.cell(0, col_idx), header, 11.5, WHITE, fill=SJTU_RED, bold=True, center=True)

    for row_idx, row in enumerate(rows[: rows_count - 1], start=1):
        for col_idx in range(cols_count):
            value = row[col_idx] if col_idx < len(row) else ""
            fill = WARM_GRAY_1 if row_idx % 2 == 0 else WHITE
            _style_table_cell(
                table.cell(row_idx, col_idx),
                value,
                10.6,
                TEXT_PRIMARY if col_idx == 0 else TEXT_BODY,
                fill=fill,
                bold=col_idx == 0,
                center=False,
            )


def _style_table_cell(
    cell,
    text: str,
    size: float,
    color: str,
    fill: str,
    bold: bool = False,
    center: bool = False,
) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(fill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(6)
    cell.margin_right = Pt(6)
    cell.margin_top = Pt(4)
    cell.margin_bottom = Pt(4)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = _trim(text, 44)
    _style_paragraph(p, size, color, bold=bold, align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT)


def _fallback_table(items: list[str]) -> tuple[list[str], list[list[str]]]:
    headers = ["项目", "说明"]
    rows: list[list[str]] = []
    for item in items[:5]:
        label, value = _split_key_value(item)
        rows.append([label, value or item])
    return headers, rows


def _render_stat_callouts(slide, items: list[str], support: str | None) -> None:
    lead, rest = _split_lead(items)
    _lead_callout(slide, lead or "用几个关键指标先建立判断。", compact=True)
    metrics = [_metric_from_item(item) for item in (rest or items)[:4]]
    if len(metrics) <= 2:
        positions = [
            (0.92, 2.25, 5.55, 2.45),
            (6.82, 2.25, 5.55, 2.45),
        ]
    elif len(metrics) == 3:
        positions = [
            (0.92, 2.25, 3.64, 2.55),
            (4.84, 2.25, 3.64, 2.55),
            (8.76, 2.25, 3.64, 2.55),
        ]
    else:
        positions = [
            (0.92, 2.18, 5.55, 1.75),
            (6.82, 2.18, 5.55, 1.75),
            (0.92, 4.22, 5.55, 1.75),
            (6.82, 4.22, 5.55, 1.75),
        ]

    for idx, metric in enumerate(metrics, start=1):
        left, top, width, height = positions[idx - 1]
        _stat_card(slide, left, top, width, height, *metric, accent=SJTU_RED if idx % 2 else GOLD)

    if support:
        _insight(slide, f"支撑：{support}", top=6.08)


def _render_metric_strip(slide, items: list[str], support: str | None) -> None:
    lead, rest = _split_lead(items)
    _lead_callout(slide, lead or "先用关键数字建立汇报判断。", compact=True)
    metrics = [_metric_from_item(item) for item in (rest or items)[:4]]
    metric_w = 2.72
    for idx, metric in enumerate(metrics, start=1):
        left = 0.92 + (idx - 1) * 2.94
        _strip_metric_card(slide, left, 2.18, metric_w, 1.26, *metric)

    detail_items = (rest or items)[4:6]
    if detail_items:
        positions = [(0.92, 3.78, 5.55, 1.86), (6.82, 3.78, 5.55, 1.86)]
        for idx, item in enumerate(detail_items[:2], start=1):
            left, top, width, height = positions[idx - 1]
            _workstream_panel(slide, left, top, width, height, idx, item, active=idx == 1)
    elif support:
        _insight(slide, support, top=5.65)


def _strip_metric_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    number: str,
    unit: str,
    desc: str,
) -> None:
    _round_rect(slide, left, top, width, height, "FDF1F1", radius=0.035)
    num_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.22), Inches(width - 0.36), Inches(0.46))
    p = num_box.text_frame.paragraphs[0]
    p.text = f"{number}{unit}"
    _style_paragraph(p, 28 if len(number) <= 5 else 23, SJTU_RED, bold=True)
    label_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.76), Inches(width - 0.4), Inches(0.24))
    p = label_box.text_frame.paragraphs[0]
    p.text = _trim(desc or label, 28)
    _style_paragraph(p, 10.6, TEXT_BODY)


def _workstream_panel(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    idx: int,
    text: str,
    active: bool = False,
) -> None:
    title, body = _split_key_value(text)
    accent = SJTU_RED if active else WARM_GRAY_3
    fill = WHITE if active else WARM_GRAY_1
    border = SJTU_RED if active else WARM_GRAY_2
    _round_rect(slide, left, top, width, height, fill, border=border, radius=0.035, shadow=True)
    _rect(slide, left, top, 0.08, height, accent)
    _number_badge(slide, left + 0.28, top + 0.26, str(idx), size=0.48, color=SJTU_RED if active else TEXT_CAPTION)
    title_box = slide.shapes.add_textbox(Inches(left + 0.92), Inches(top + 0.24), Inches(width - 1.2), Inches(0.32))
    p = title_box.text_frame.paragraphs[0]
    p.text = _trim(title, 26)
    _style_paragraph(p, 14.5, SJTU_RED if active else TEXT_PRIMARY, bold=True)
    _draw_check_points(slide, left + 0.46, top + 0.78, width - 0.76, height - 0.9, _split_card_points(body), compact=True)


def _render_milestone_timeline(slide, items: list[str]) -> None:
    lead, rest = _split_lead(items)
    if lead:
        _lead_callout(slide, lead, compact=True)
    milestones = rest[:4] or items[:4]
    if not milestones:
        milestones = ["阶段一：明确目标；完成启动", "阶段二：交付原型；验证效果", "阶段三：扩大应用；沉淀经验"]

    axis_y = 2.34 if lead else 1.72
    _rect(slide, 1.2, axis_y, 10.88, 0.035, SJTU_RED)
    count = min(len(milestones), 4)
    col_w = 11.48 / max(count, 1)
    for idx, item in enumerate(milestones[:4], start=1):
        x = 0.92 + (idx - 1) * col_w
        center = x + col_w / 2
        _number_badge(slide, center - 0.15, axis_y - 0.15, str(idx), size=0.34, color=SJTU_RED)
        label, body = _split_key_value(item)
        title_box = slide.shapes.add_textbox(Inches(x + 0.08), Inches(axis_y + 0.35), Inches(col_w - 0.16), Inches(0.34))
        p = title_box.text_frame.paragraphs[0]
        p.text = _trim(label, 14)
        _style_paragraph(p, 13.4, SJTU_RED, bold=True, align=PP_ALIGN.CENTER)
        _timeline_stage_card(slide, x + 0.08, axis_y + 0.86, col_w - 0.18, 1.24, body or item)
        _image_placeholder(slide, x + 0.08, axis_y + 2.36, col_w - 0.18, 1.48, "系统界面截图")


def _timeline_stage_card(slide, left: float, top: float, width: float, height: float, body: str) -> None:
    _round_rect(slide, left, top, width, height, "FDF1F1", radius=0.035)
    _draw_check_points(slide, left + 0.22, top + 0.18, width - 0.44, height - 0.24, _split_card_points(body), compact=True)


def _render_project_showcase(slide, items: list[str]) -> None:
    lead, rest = _split_lead(items)
    if lead:
        _lead_callout(slide, lead, compact=True)
    projects = rest[:4] or items[:4]
    top = 2.12 if lead else 1.5
    card_w = 2.73
    gap = 0.24
    for idx, item in enumerate(projects[:4], start=1):
        left = 0.92 + (idx - 1) * (card_w + gap)
        _project_showcase_card(slide, left, top, card_w, 2.28, item)
        title, _body = _split_key_value(item)
        _image_placeholder(slide, left, top + 2.52, card_w, 1.48, title or f"项目 {idx}")


def _project_showcase_card(slide, left: float, top: float, width: float, height: float, text: str) -> None:
    title, body = _split_key_value(text)
    if not body:
        title, body = _trim(text, 18), text
    _rect(slide, left, top, width, 0.78, SJTU_RED)
    title_box = slide.shapes.add_textbox(Inches(left + 0.28), Inches(top + 0.18), Inches(width - 0.5), Inches(0.36))
    p = title_box.text_frame.paragraphs[0]
    p.text = _trim(title, 16)
    _style_paragraph(p, 13.2, WHITE, bold=True)
    _rect(slide, left, top + 0.78, width, height - 0.78, WARM_GRAY_1)
    _draw_check_points(slide, left + 0.28, top + 1.02, width - 0.5, height - 1.08, _split_card_points(body), compact=True)


def _render_media_explain(slide, items: list[str], support: str | None) -> None:
    lead, rest = _split_lead(items)
    if lead:
        _lead_callout(slide, lead, compact=True)
    top = 2.05 if lead else 1.42
    _image_placeholder(slide, 0.92, top, 6.15, 3.95, support or "系统界面截图")
    explain = rest[:4] or items[:4]
    right_x = 7.45
    section_top = top + 0.12
    first = explain[0] if explain else "定位：说明这个产品、平台或方法解决什么问题"
    title, body = _split_key_value(first)
    _section_heading(slide, right_x, section_top, title or "定位")
    body_box = slide.shapes.add_textbox(Inches(right_x), Inches(section_top + 0.48), Inches(4.72), Inches(0.88))
    body_box.text_frame.word_wrap = True
    p = body_box.text_frame.paragraphs[0]
    p.text = _trim(body or first, 98)
    _style_paragraph(p, 12.8, TEXT_BODY)
    p.line_spacing = 1.12

    stack_items = explain[1:4]
    if stack_items:
        _section_heading(slide, right_x, section_top + 1.58, "结构")
        for idx, item in enumerate(stack_items[:3]):
            label, value = _split_key_value(item)
            y = section_top + 2.02 + idx * 0.48
            color = SJTU_RED if idx == 0 else "D74444" if idx == 1 else "8C8C8C"
            shape = _round_rect(slide, right_x + 0.58 * idx, y, 4.2 - 0.58 * idx, 0.34, color, radius=0.04)
            _shape_text(shape, _trim(label if value else item, 24), 10.5, WHITE, bold=True, align=PP_ALIGN.CENTER)
    point_items = explain[4:]
    if point_items:
        _draw_check_points(slide, right_x, top + 3.56, 4.72, 0.82, point_items[:3], compact=True)


def _metric_from_item(item: str) -> tuple[str, str, str, str]:
    label, value = _split_key_value(item)
    if not value:
        return "指标", _trim(item, 16), "", ""
    parts = [part.strip(" 。") for part in re.split(r"\s*/\s*|\s*[；;]\s*", value) if part.strip(" 。")]
    number = parts[0] if parts else value
    unit = ""
    desc_parts = parts[1:] if len(parts) > 1 else []
    number_match = re.match(r"^([~≈]?\d+(?:\.\d+)?)([%％A-Za-z一-龥]*)$", number)
    if number_match:
        number = number_match.group(1)
        unit = number_match.group(2)
    return _trim(label, 16), _trim(number, 18), _trim(unit, 8), _trim("；".join(desc_parts), 64)


def _stat_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    number: str,
    unit: str,
    desc: str,
    accent: str,
) -> None:
    _round_rect(slide, left, top, width, height, WHITE, border=WARM_GRAY_2, radius=0.035, shadow=True)
    _rect(slide, left, top, width, 0.05, accent)
    label_box = slide.shapes.add_textbox(Inches(left + 0.24), Inches(top + 0.18), Inches(width - 0.48), Inches(0.28))
    p = label_box.text_frame.paragraphs[0]
    p.text = label
    _style_paragraph(p, 10.8, TEXT_CAPTION, bold=True)

    number_box = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.54), Inches(width - 0.44), Inches(0.72))
    p = number_box.text_frame.paragraphs[0]
    p.text = f"{number}{unit}"
    _style_paragraph(p, 27 if width < 4 else 32, accent, bold=True)

    if desc:
        desc_box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 1.28), Inches(width - 0.5), Inches(height - 1.46))
        desc_box.text_frame.word_wrap = True
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        _style_paragraph(p, 10.8 if width < 4 else 11.5, TEXT_BODY)
        p.line_spacing = 1.12


def _render_quote_block_slide(slide, items: list[str], support: str | None) -> None:
    lead, rest = _split_lead(items)
    if lead:
        _lead_callout(slide, lead, compact=True)
    quote = rest[0] if rest else lead or items[0]
    quote_title, quote_body = _split_key_value(quote)
    if not quote_body:
        quote_title, quote_body = "核心判断", quote

    _quote_block_wide(slide, 0.92, 2.08, 11.55, 2.15, quote_body, quote_title, support)

    supporting_items = rest[1:3] if rest else items[1:3]
    for idx, item in enumerate(supporting_items, start=1):
        left = 0.92 if idx == 1 else 6.82
        _content_card(slide, left, 4.75, 5.55, 1.26, idx, item)


def _quote_block_wide(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    label: str,
    source: str | None,
) -> None:
    _round_rect(slide, left, top, width, height, GOLD_PALE, border=GOLD, radius=0.035, shadow=True)
    _rect(slide, left, top, 0.06, height, GOLD)
    mark = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.18), Inches(0.55), Inches(0.58))
    p = mark.text_frame.paragraphs[0]
    p.text = '"'
    _style_paragraph(p, 30, GOLD, bold=True, align=PP_ALIGN.CENTER)

    label_box = slide.shapes.add_textbox(Inches(left + 0.86), Inches(top + 0.28), Inches(2.2), Inches(0.28))
    p = label_box.text_frame.paragraphs[0]
    p.text = _trim(label, 18)
    _style_paragraph(p, 10.5, SJTU_RED, bold=True)

    body_box = slide.shapes.add_textbox(Inches(left + 0.86), Inches(top + 0.68), Inches(width - 1.35), Inches(0.72))
    body_box.text_frame.word_wrap = True
    p = body_box.text_frame.paragraphs[0]
    p.text = _trim(text, 92)
    _style_paragraph(p, 19 if len(text) < 52 else 15.5, TEXT_PRIMARY, bold=True)
    p.line_spacing = 1.08

    if source:
        source_box = slide.shapes.add_textbox(Inches(left + 0.86), Inches(top + 1.58), Inches(width - 1.35), Inches(0.28))
        p = source_box.text_frame.paragraphs[0]
        p.text = f"支撑：{_trim(source, 80)}"
        _style_paragraph(p, 10.2, TEXT_CAPTION)


def _render_card_grid(slide, items: list[str]) -> None:
    lead, cards = _split_lead(items)
    if not lead and items:
        lead, cards = items[0], items[1:]
    _lead_callout(slide, lead)
    cards = cards[:4] or items[:4]
    positions = [
        (0.92, 2.42),
        (6.82, 2.42),
        (0.92, 4.45),
        (6.82, 4.45),
    ]
    for idx, item in enumerate(cards[:4], start=1):
        left, top = positions[idx - 1]
        _content_card(slide, left, top, 5.55, 1.55, idx, item)


def _render_rich_card_grid(slide, items: list[str]) -> None:
    lead, cards = _split_lead(items)
    if not lead and items:
        lead, cards = items[0], items[1:]
    _lead_callout(slide, lead, compact=True)
    cards = cards[:4] or items[:4]
    count = len(cards)
    if count <= 2:
        positions = [
            (0.92, 2.26, 5.55, 3.28),
            (6.82, 2.26, 5.55, 3.28),
        ]
    elif count == 3:
        positions = [
            (0.92, 2.16, 3.64, 3.74),
            (4.84, 2.16, 3.64, 3.74),
            (8.76, 2.16, 3.64, 3.74),
        ]
    else:
        positions = [
            (0.92, 2.16, 5.55, 1.86),
            (6.82, 2.16, 5.55, 1.86),
            (0.92, 4.36, 5.55, 1.86),
            (6.82, 4.36, 5.55, 1.86),
        ]
    for idx, item in enumerate(cards, start=1):
        left, top, width, height = positions[idx - 1]
        _rich_content_card(slide, left, top, width, height, idx, item)


def _render_fact_grid(slide, items: list[str]) -> None:
    lead, rest = _split_lead(items)
    key_values = [_split_key_value(item) for item in rest]
    if not key_values:
        key_values = [_split_key_value(item) for item in items]
        lead = ""
    if lead:
        _lead_callout(slide, lead, compact=True)
        top_start = 2.15
    else:
        top_start = 1.28
    for idx, (label, value) in enumerate(key_values[:6], start=1):
        col = (idx - 1) % 2
        row = (idx - 1) // 2
        left = 0.92 + col * 5.9
        top = top_start + row * 1.45
        _fact_card(slide, left, top, 5.55, 1.08, label or f"要点 {idx}", value)


def _render_timeline(slide, items: list[str]) -> None:
    lead, rest = _split_lead(items)
    _lead_callout(slide, lead, compact=True)
    timeline_items = rest[:4] if rest else items[:4]
    axis_x = 1.22
    top0 = 2.08
    _rect(slide, axis_x + 0.13, top0 + 0.16, 0.025, 3.7, GOLD)
    for idx, item in enumerate(timeline_items, start=1):
        top = top0 + (idx - 1) * 0.95
        _number_badge(slide, axis_x, top, str(idx), size=0.32, color=SJTU_RED if idx % 2 else GOLD)
        date, text = _split_date(item)
        if date:
            date_box = slide.shapes.add_textbox(Inches(1.72), Inches(top - 0.03), Inches(1.55), Inches(0.28))
            p = date_box.text_frame.paragraphs[0]
            p.text = date
            _style_paragraph(p, 11, SJTU_RED, bold=True)
            text_left = 3.2
            text_w = 5.45
        else:
            text_left = 1.72
            text_w = 6.95
        text_box = slide.shapes.add_textbox(Inches(text_left), Inches(top - 0.04), Inches(text_w), Inches(0.45))
        p = text_box.text_frame.paragraphs[0]
        p.text = _trim(text, 78)
        _style_paragraph(p, 12.5, TEXT_BODY)

    _quote_panel(
        slide,
        9.1,
        2.05,
        3.25,
        3.65,
        "这个速度，值得认真对待。",
        "从“能聊天”到“能独立做数值计算科研”",
    )


def _render_process(slide, items: list[str]) -> None:
    lead, rest = _split_lead(items)
    _lead_callout(slide, lead, compact=True)
    formula_items = [item for item in rest if _is_formula_item(item)]
    non_formula_items = [item for item in rest if item not in formula_items]
    process = non_formula_items[:3] if formula_items else non_formula_items[:4]
    process = process or items[:4]
    left = 0.9
    top = 2.38
    width = 11.55
    gap = width / max(len(process), 1)
    colors = [SJTU_RED, GOLD, BROWN, TEXT_PRIMARY]
    for idx, item in enumerate(process, start=1):
        x = left + (idx - 1) * gap
        card_w = max(gap - 0.22, 2.25)
        accent = colors[(idx - 1) % len(colors)]
        _process_card(slide, x, top, card_w, 2.18, idx, item, accent)

    if formula_items:
        _formula_panel(slide, _normalize_formula_text(formula_items[0]), top=5.08)
        return

    if len(items) > len(process) + 1:
        _insight(slide, items[-1], top=5.75)


def _render_concept_diagram(slide, items: list[str], support: str | None) -> None:
    lead, rest = _split_lead(items)
    _lead_callout(slide, lead, compact=True)
    nodes = _concept_nodes(rest, support)

    _concept_connector(slide, 4.17, 3.48, 5.05, 3.48)
    _concept_connector(slide, 8.3, 3.48, 9.05, 3.48)
    _concept_connector(slide, 6.67, 4.66, 6.67, 5.02)

    positions = [
        (0.92, 2.52, 3.25, 1.92, SJTU_RED, False),
        (5.05, 2.34, 3.25, 2.28, GOLD, True),
        (9.05, 2.52, 3.25, 1.92, BROWN, False),
        (4.35, 5.02, 4.65, 0.82, SJTU_RED, False),
    ]
    for (title, body), (left, top, width, height, accent, emphasis) in zip(nodes, positions, strict=False):
        _concept_node(slide, left, top, width, height, title, body, accent, emphasis=emphasis)


def _concept_nodes(items: list[str], support: str | None) -> list[tuple[str, str]]:
    nodes: list[tuple[str, str]] = []
    for item in items:
        if _is_formula_item(item):
            continue
        title, body = _split_key_value(item)
        if not body:
            title, body = _trim(item, 18), item
        nodes.append((_trim(title, 14), _normalize_formula_text(body)))
    fallbacks = [
        ("输入", "先明确对象、材料和可观察信号。"),
        ("模型", "把输入转成判断、预测或行动。"),
        ("输出", "得到可检查的结果，再和目标比较。"),
        ("反馈", support or "根据误差或评价结果，回到前一步继续改进。"),
    ]
    for fallback in fallbacks:
        if len(nodes) >= 4:
            break
        nodes.append(fallback)
    return nodes[:4]


def _concept_connector(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    if abs(x2 - x1) >= abs(y2 - y1):
        left = min(x1, x2)
        _rect(slide, left, y1, max(abs(x2 - x1), 0.04), 0.035, GOLD)
    else:
        top = min(y1, y2)
        _rect(slide, x1, top, 0.035, max(abs(y2 - y1), 0.04), GOLD)


def _concept_node(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: str,
    accent: str,
    emphasis: bool = False,
) -> None:
    fill = GOLD_PALE if emphasis else WHITE
    border = accent if emphasis else WARM_GRAY_2
    _round_rect(slide, left, top, width, height, fill, border=border, radius=0.035, shadow=True)
    _rect(slide, left, top, width, 0.05, accent)
    if height <= 1.1:
        label_box = slide.shapes.add_textbox(Inches(left + 0.28), Inches(top + 0.16), Inches(1.15), Inches(0.28))
        p = label_box.text_frame.paragraphs[0]
        p.text = _trim(title, 8)
        _style_paragraph(p, 12, accent, bold=True)
        body_box = slide.shapes.add_textbox(Inches(left + 1.48), Inches(top + 0.14), Inches(width - 1.75), Inches(0.46))
        body_box.text_frame.word_wrap = True
        body_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = body_box.text_frame.paragraphs[0]
        p.text = _trim(body, 70)
        _style_paragraph(p, 10.6, TEXT_BODY, font_name=MATH_FONT if _looks_like_formula_text(body) else FONT)
        return

    label_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.25), Inches(width - 0.6), Inches(0.38))
    label_box.text_frame.word_wrap = True
    label_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_box.text_frame.paragraphs[0]
    p.text = _trim(title, 16)
    _style_paragraph(p, 14 if emphasis else 13.2, TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    _draw_card_points(
        slide,
        left + 0.3,
        top + 0.83,
        width - 0.6,
        height - 1.0,
        _split_card_points(body),
        compact=True,
    )


def _render_example_walkthrough(slide, items: list[str]) -> None:
    lead, rest = _split_lead(items)
    _lead_callout(slide, lead, compact=True)
    formula_items = [item for item in rest if _is_formula_item(item)]
    non_formula_items = [item for item in rest if item not in formula_items]
    cards = non_formula_items[:3] or rest[:3] or items[1:4]
    if not cards:
        cards = ["输入：先明确样本和标签", "模型：写出预测关系", "验证：检查是否能迁移"]

    left = 0.9
    top = 2.34
    width = 11.55
    gap = 0.32
    card_w = (width - gap * 2) / 3
    colors = [SJTU_RED, GOLD, BROWN]
    _example_relation_line(slide, left, top + 1.12, card_w, gap, len(cards[:3]))
    for idx, item in enumerate(cards[:3], start=1):
        x = left + (idx - 1) * (card_w + gap)
        _example_walkthrough_card(slide, x, top, card_w, 2.25, idx, item, colors[idx - 1])

    if formula_items:
        _example_formula_board(slide, _normalize_formula_text(formula_items[0]), top=5.12)
    elif len(non_formula_items) > 3:
        _insight(slide, non_formula_items[3], top=5.38)


def _example_relation_line(slide, left: float, top: float, card_w: float, gap: float, count: int) -> None:
    if count < 2:
        return
    for idx in range(count - 1):
        x = left + (idx + 1) * card_w + idx * gap + 0.03
        _rect(slide, x, top, max(0.05, gap - 0.06), 0.035, GOLD)


def _lead_callout(slide, text: str, compact: bool = False) -> None:
    text = text or "核心结论"
    h = 0.88 if compact else 1.05
    card_left, card_top, card_width = 0.9, 1.16, 11.55
    _round_rect(slide, card_left, card_top, card_width, h, WHITE, border=WARM_GRAY_2, radius=0.04, shadow=True)
    _rect(slide, card_left, card_top, 0.06, h, GOLD)
    box = slide.shapes.add_textbox(Inches(1.17), Inches(card_top), Inches(10.8), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    p = frame.paragraphs[0]
    p.text = _trim(text, 104)
    _style_paragraph(p, 17.6 if len(text) < 58 else 15.2, TEXT_PRIMARY, bold=True)
    p.line_spacing = 1.08


def _content_card(slide, left: float, top: float, width: float, height: float, idx: int, text: str) -> None:
    _round_rect(slide, left, top, width, height, WHITE, border=WARM_GRAY_2, radius=0.035, shadow=True)
    _rect(slide, left, top, width, 0.04, SJTU_RED if idx % 2 else GOLD)
    _number_badge(slide, left + 0.24, top + 0.23, str(idx), size=0.38, color=SJTU_RED if idx % 2 else GOLD)
    title, body = _split_key_value(text)
    title_box = slide.shapes.add_textbox(Inches(left + 0.78), Inches(top + 0.18), Inches(width - 1.08), Inches(0.32))
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_box.text_frame.paragraphs[0]
    p.text = title if body else _trim(text, 36)
    _style_paragraph(p, 14, TEXT_PRIMARY, bold=True)
    if body:
        body_text = body
        body_top = top + 0.58
    else:
        body_text = text
        body_top = top + 0.56
    body_box = slide.shapes.add_textbox(Inches(left + 0.34), Inches(body_top), Inches(width - 0.62), Inches(height - 0.72))
    body_box.text_frame.word_wrap = True
    body_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = body_box.text_frame.paragraphs[0]
    body_text = _normalize_formula_text(body_text)
    p.text = _trim(body_text, 96)
    if _looks_like_formula_text(body_text):
        _style_paragraph(p, 14.5 if len(body_text) <= 58 else 12.8, TEXT_BODY, font_name=MATH_FONT)
    else:
        _style_paragraph(p, 12.2 if len(body_text) > 52 else 13, TEXT_BODY)
    p.line_spacing = 1.15


def _rich_content_card(slide, left: float, top: float, width: float, height: float, idx: int, text: str) -> None:
    accent = SJTU_RED if idx % 2 else GOLD
    _round_rect(slide, left, top, width, height, WHITE, border=WARM_GRAY_2, radius=0.035, shadow=True)
    _rect(slide, left, top, width, 0.05, accent)
    title, body = _split_key_value(text)
    if not body:
        title, body = _trim(text, 22), text
    label_box = slide.shapes.add_textbox(Inches(left + 0.28), Inches(top + 0.25), Inches(width - 0.56), Inches(0.42))
    label_box.text_frame.word_wrap = True
    label_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_box.text_frame.paragraphs[0]
    p.text = _trim(title, 18 if width < 4 else 24)
    _style_paragraph(p, 14 if width >= 5 else 13, TEXT_PRIMARY, bold=True)
    _draw_card_points(
        slide,
        left + 0.3,
        top + 0.86,
        width - 0.6,
        height - 1.08,
        _split_card_points(body),
        compact=width < 4.2,
    )


def _example_walkthrough_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    idx: int,
    text: str,
    accent: str,
) -> None:
    _round_rect(slide, left, top, width, height, WHITE, border=WARM_GRAY_2, radius=0.035, shadow=True)
    _rect(slide, left, top, width, 0.05, accent)
    _number_badge(slide, left + 0.2, top + 0.22, str(idx), size=0.36, color=accent)
    title, body = _split_process_text(text)

    title_box = slide.shapes.add_textbox(Inches(left + 0.68), Inches(top + 0.18), Inches(width - 0.92), Inches(0.45))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_box.text_frame.paragraphs[0]
    p.text = _trim(title, 18)
    _style_paragraph(p, 13.4 if len(title) <= 10 else 12.7, TEXT_PRIMARY, bold=True)

    points = _split_card_points(body)
    _draw_card_points(
        slide,
        left + 0.28,
        top + 0.82,
        width - 0.56,
        height - 0.98,
        points,
        compact=True,
    )


def _fact_card(slide, left: float, top: float, width: float, height: float, label: str, value: str) -> None:
    _round_rect(slide, left, top, width, height, WHITE, border=WARM_GRAY_2, radius=0.035, shadow=True)
    _rect(slide, left, top, 0.05, height, GOLD)
    label_box = slide.shapes.add_textbox(Inches(left + 0.23), Inches(top + 0.16), Inches(1.32), Inches(0.26))
    label_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_box.text_frame.paragraphs[0]
    p.text = _trim(label, 10)
    _style_paragraph(p, 10.5, SJTU_RED, bold=True)
    value_box = slide.shapes.add_textbox(Inches(left + 1.45), Inches(top + 0.14), Inches(width - 1.7), Inches(height - 0.2))
    value_box.text_frame.word_wrap = True
    value_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = value_box.text_frame.paragraphs[0]
    value = _normalize_formula_text(value)
    p.text = _trim(value, 92)
    if _looks_like_formula_text(value):
        _style_paragraph(p, 12.6 if len(value) < 58 else 10.8, TEXT_BODY, font_name=MATH_FONT)
    else:
        _style_paragraph(p, 12.2 if len(value) < 54 else 10.8, TEXT_BODY)
    p.line_spacing = 1.12


def _process_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    idx: int,
    text: str,
    accent: str,
) -> None:
    _round_rect(slide, left, top, width, height, WHITE, border=WARM_GRAY_2, radius=0.035, shadow=True)
    _rect(slide, left, top, width, 0.05, accent)
    _number_badge(slide, left + 0.2, top + 0.22, str(idx), size=0.36, color=accent)
    title, body = _split_process_text(text)
    title_box = slide.shapes.add_textbox(Inches(left + 0.65), Inches(top + 0.18), Inches(width - 0.85), Inches(0.48))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_box.text_frame.paragraphs[0]
    p.text = _trim(title, 20)
    _style_paragraph(p, 12.5 if len(title) > 12 else 13.4, TEXT_PRIMARY, bold=True)
    p.line_spacing = 1.0
    body = _normalize_formula_text(body)
    points = _split_card_points(body)
    if len(points) > 1:
        _draw_card_points(
            slide,
            left + 0.24,
            top + 0.8,
            width - 0.48,
            height - 0.94,
            points,
            compact=True,
        )
    else:
        body_box = slide.shapes.add_textbox(Inches(left + 0.24), Inches(top + 0.82), Inches(width - 0.48), Inches(0.62))
        body_box.text_frame.word_wrap = True
        body_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = body_box.text_frame.paragraphs[0]
        p.text = _trim(body, 72)
        if _looks_like_formula_text(body):
            _style_paragraph(p, 12 if len(body) <= 44 else 11.2, TEXT_BODY, font_name=MATH_FONT)
        else:
            _style_paragraph(p, 11 if len(body) > 34 else 11.6, TEXT_BODY)
        p.alignment = PP_ALIGN.CENTER


def _draw_card_points(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    points: list[str],
    compact: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    max_points = 3 if compact else 4
    point_font_size = 11.4 if compact and width < 2.35 else 11.8 if compact else 12.1
    trim_chars = 32 if compact and width < 2.35 else 36 if compact else 54
    for idx, point in enumerate(points[:max_points]):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        text = _normalize_formula_text(point)
        p.text = f"▪ {_trim(text, trim_chars)}"
        font_name = MATH_FONT if _looks_like_formula_text(text) else FONT
        _style_paragraph(
            p,
            point_font_size,
            TEXT_BODY,
            font_name=font_name,
        )
        p.line_spacing = 1.08
        p.space_after = Pt(6 if compact else 8)


def _draw_check_points(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    points: list[str],
    compact: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    max_points = 3 if compact else 4
    trim_chars = 28 if compact and width < 3 else 42 if compact else 58
    for idx, point in enumerate((points or ["请补充要点"])[:max_points]):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        text = _normalize_formula_text(point)
        p.text = ""
        check = p.add_run()
        check.text = "✓ "
        check.font.name = FONT
        check.font.size = Pt(11.2 if compact else 12.2)
        check.font.bold = True
        check.font.color.rgb = rgb(SJTU_RED)
        _set_ea(check)
        run = p.add_run()
        run.text = _trim(text, trim_chars)
        run.font.name = FONT
        run.font.size = Pt(11.2 if compact else 12.2)
        run.font.color.rgb = rgb(TEXT_BODY)
        _set_ea(run)
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.08
        p.space_after = Pt(5 if compact else 7)


def _section_heading(slide, left: float, top: float, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(4.7), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = _trim(text, 18)
    _style_paragraph(p, 13, SJTU_RED, bold=True)


def _image_placeholder(slide, left: float, top: float, width: float, height: float, caption: str) -> None:
    _round_rect(slide, left, top, width, height, "F1F4F7", border="C8D2DF", radius=0.035)
    icon = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + width / 2 - 0.26),
        Inches(top + height / 2 - 0.18),
        Inches(0.52),
        Inches(0.36),
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = rgb("9AA7B5")
    icon.line.fill.background()
    icon.shadow.inherit = False
    _shape_text(icon, "图", 10.5, WHITE, bold=True, align=PP_ALIGN.CENTER)
    label = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + height - 0.36), Inches(width - 0.4), Inches(0.22))
    p = label.text_frame.paragraphs[0]
    p.text = f"建议配图：{_trim(caption, 28)}"
    _style_paragraph(p, 8.8, TEXT_CAPTION, align=PP_ALIGN.CENTER)


def _quote_panel(slide, left: float, top: float, width: float, height: float, title: str, body: str) -> None:
    _round_rect(slide, left, top, width, height, GOLD_PALE, border=GOLD, radius=0.035, shadow=True)
    _rect(slide, left, top, width, 0.05, SJTU_RED)
    title_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.38), Inches(width - 0.6), Inches(0.8))
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    _style_paragraph(p, 20, SJTU_DEEP, bold=True, align=PP_ALIGN.CENTER)
    body_box = slide.shapes.add_textbox(Inches(left + 0.35), Inches(top + 1.58), Inches(width - 0.7), Inches(1.05))
    body_box.text_frame.word_wrap = True
    p = body_box.text_frame.paragraphs[0]
    p.text = body
    _style_paragraph(p, 12, TEXT_BODY, align=PP_ALIGN.CENTER)


def _insight(slide, text: str, top: float = 6.05) -> None:
    _round_rect(slide, 0.9, top, 11.55, 0.5, WARM_GRAY_1, border=WARM_GRAY_2, radius=0.045)
    _rect(slide, 0.9, top, 0.04, 0.5, GOLD)
    box = slide.shapes.add_textbox(Inches(1.12), Inches(top + 0.08), Inches(11.0), Inches(0.34))
    p = box.text_frame.paragraphs[0]
    p.text = _trim(text, 98)
    _style_paragraph(p, 11.4, TEXT_BODY)


def _formula_panel(slide, text: str, top: float = 4.88) -> None:
    title, formula = _split_key_value(text)
    if not formula:
        title, formula = "公式", text
    _round_rect(slide, 0.9, top, 11.55, 0.86, GOLD_PALE, border=GOLD, radius=0.035)
    _rect(slide, 0.9, top, 0.05, 0.86, SJTU_RED)
    label_box = slide.shapes.add_textbox(Inches(1.16), Inches(top + 0.14), Inches(1.1), Inches(0.24))
    p = label_box.text_frame.paragraphs[0]
    p.text = _trim(title, 10)
    _style_paragraph(p, 10.5, SJTU_RED, bold=True)
    formula_box = slide.shapes.add_textbox(Inches(2.1), Inches(top + 0.16), Inches(9.85), Inches(0.5))
    formula_box.text_frame.word_wrap = True
    p = formula_box.text_frame.paragraphs[0]
    p.text = _trim(_normalize_formula_text(formula), 118)
    _style_paragraph(p, 15 if len(formula) <= 70 else 12.5, TEXT_PRIMARY, font_name=MATH_FONT)


def _example_formula_board(slide, text: str, top: float = 5.08) -> None:
    title, formula = _split_key_value(text)
    if not formula:
        title, formula = "板书", text
    label = "板书" if title in {"公式", "损失函数", "目标函数"} else title
    formula = _normalize_formula_text(formula)
    _round_rect(slide, 0.9, top, 11.55, 0.98, GOLD_PALE, border=GOLD, radius=0.035, shadow=True)
    _rect(slide, 0.9, top, 0.06, 0.98, SJTU_RED)

    label_box = slide.shapes.add_textbox(Inches(1.18), Inches(top + 0.16), Inches(1.35), Inches(0.32))
    p = label_box.text_frame.paragraphs[0]
    p.text = _trim(label, 8)
    _style_paragraph(p, 13, SJTU_RED, bold=True)

    formula_box = slide.shapes.add_textbox(Inches(2.55), Inches(top + 0.15), Inches(9.0), Inches(0.38))
    formula_box.text_frame.word_wrap = True
    p = formula_box.text_frame.paragraphs[0]
    p.text = _trim(formula, 92)
    _style_paragraph(p, 16 if len(formula) <= 62 else 13.6, TEXT_PRIMARY, font_name=MATH_FONT)

    note_box = slide.shapes.add_textbox(Inches(2.55), Inches(top + 0.58), Inches(9.3), Inches(0.24))
    p = note_box.text_frame.paragraphs[0]
    p.text = "把误差变小，就是让模型参数更适合这批训练样本。"
    _style_paragraph(p, 10.8, TEXT_BODY)


def render_footer(slide, page_num: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.85), Inches(6.92), Inches(0.85), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.text = f"◆  {page_num:02d}"
    _style_paragraph(p, 9, TEXT_CAPTION)
    p.alignment = PP_ALIGN.RIGHT


def _fill_background(slide, color: str) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(color)


def _rect(slide, left: float, top: float, width: float, height: float, fill: str, border: str | None = None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if border:
        shape.line.color.rgb = rgb(border)
        shape.line.width = Pt(0.6)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _round_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: str,
    border: str | None = None,
    radius: float = 0.045,
    shadow: bool = False,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if border:
        shape.line.color.rgb = rgb(border)
        shape.line.width = Pt(0.55)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = radius
    except (IndexError, TypeError):
        pass
    if shadow:
        _add_shadow(shape)
    return shape


def _number_badge(
    slide,
    left: float,
    top: float,
    text: str,
    size: float = 0.36,
    color: str = SJTU_RED,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    _shape_text(shape, text, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)
    return shape


def _decorative_rule(slide, left: float, top: float, width: float) -> None:
    _rect(slide, left, top, width, 0.018, GOLD)
    diamond = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(left + width + 0.1),
        Inches(top - 0.035),
        Inches(0.09),
        Inches(0.09),
    )
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = rgb(GOLD)
    diamond.line.fill.background()


def _section_label(slide, label: str, subtitle: str) -> None:
    pill = _round_rect(slide, 0.9, 1.1, 1.35, 0.34, SJTU_RED, radius=0.14)
    _shape_text(pill, label, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)
    box = slide.shapes.add_textbox(Inches(2.45), Inches(1.08), Inches(6.8), Inches(0.38))
    p = box.text_frame.paragraphs[0]
    p.text = subtitle
    _style_paragraph(p, 12.2, TEXT_CAPTION)


def _shape_text(shape, text: str, size: float, color: str, bold: bool = False, align=PP_ALIGN.LEFT):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    p = frame.paragraphs[0]
    p.text = text
    _style_paragraph(p, size, color, bold=bold, align=align)
    return p


def _style_paragraph(
    p,
    size: float,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = FONT,
) -> None:
    p.font.name = font_name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.alignment = align
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        _set_ea(run, font_name)


def _set_ea(run, font_name: str = FONT) -> None:
    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find(qn("a:ea"))
    if ea is None:
        ea = r_pr.makeelement(qn("a:ea"), {})
        r_pr.append(ea)
    ea.set("typeface", font_name)


def _add_shadow(shape) -> None:
    sp_pr = shape._element.spPr
    effect_lst = sp_pr.find(qn("a:effectLst"))
    if effect_lst is None:
        effect_lst = sp_pr.makeelement(qn("a:effectLst"), {})
        sp_pr.append(effect_lst)
    else:
        for child in list(effect_lst):
            effect_lst.remove(child)
    outer = sp_pr.makeelement(qn("a:outerShdw"), {})
    outer.set("blurRad", "63500")
    outer.set("dist", "19050")
    outer.set("dir", "5400000")
    outer.set("rotWithShape", "0")
    color = sp_pr.makeelement(qn("a:srgbClr"), {})
    color.set("val", SJTU_DEEP)
    alpha = sp_pr.makeelement(qn("a:alpha"), {})
    alpha.set("val", "9000")
    color.append(alpha)
    outer.append(color)
    effect_lst.append(outer)


def _clean_items(items: list[str]) -> list[str]:
    cleaned: list[str] = []
    for item in items:
        text = item.strip()
        if not text or TABLE_RULE_RE.match(text):
            continue
        if text in {"| | |", "|---|---|"}:
            continue
        text = re.sub(r"\s+", " ", text)
        text = text.strip(" -•")
        text = _strip_content_label(text)
        text = _normalize_formula_text(text)
        if text and text not in cleaned:
            cleaned.append(text)
    return cleaned[:6]


def _strip_content_label(text: str) -> str:
    return re.sub(r"^(?:一句话|核心判断|核心结论|本页核心判断)\s*[：:]\s*", "", text).strip()


def _looks_like_formula_text(text: str) -> bool:
    return bool(FORMULA_RE.search(text or "")) or text.strip().startswith(("公式：", "公式:"))


def _is_formula_item(text: str) -> bool:
    normalized = _normalize_formula_text(text)
    label, value = _split_key_value(normalized)
    if label in {"公式", "目标函数", "损失函数", "更新规则"} and value:
        return True
    return normalized.startswith(("$", r"\(", r"\["))


def _normalize_formula_text(text: str) -> str:
    text = text.strip()
    if not text:
        return text
    text = text.replace("$$", "$")
    if text.startswith("$") and text.endswith("$") and len(text) > 2:
        text = text[1:-1].strip()
    text = text.replace(r"\(", "").replace(r"\)", "")
    text = text.replace(r"\[", "").replace(r"\]", "")
    text = re.sub(r"\\frac\{([^{}]+)\}\{([^{}]+)\}", r"(\1/\2)", text)
    text = re.sub(r"_\{([^{}]+)\}", r"_\1", text)
    text = re.sub(r"\^\{([^{}]+)\}", r"^\1", text)
    for latex, replacement in LATEX_REPLACEMENTS.items():
        text = text.replace(latex, replacement)
    return re.sub(r"\s+", " ", text).strip()


def _split_lead(items: list[str]) -> tuple[str, list[str]]:
    if not items:
        return "", []
    first = items[0]
    if len(items) == 1:
        return first, []
    if _split_key_value(first)[1] and not DATE_RE.search(first):
        return "", items
    return first, items[1:]


def _split_key_value(text: str) -> tuple[str, str]:
    match = KEY_VALUE_RE.match(text)
    if match:
        return match.group(1).strip(), match.group(2).strip()
    if " / " in text:
        head, _, tail = text.partition(" / ")
        if len(head) <= 16:
            return head.strip(), tail.strip()
    return _trim(text, 28), ""


def _split_date(text: str) -> tuple[str, str]:
    match = re.match(r"^([20]\d{3}[./-]\d{1,2}[./-]\d{1,2}|20\d{2}\s*年|三年前|上周)\s*[—\-:：]?\s*(.+)$", text)
    if match:
        return match.group(1), match.group(2).strip(" ，,")
    return "", text


def _looks_like_timeline(items: list[str]) -> bool:
    return sum(1 for item in items if DATE_RE.search(item)) >= 2


def _looks_like_fact_grid(items: list[str]) -> bool:
    candidates = items
    if items:
        label, value = _split_key_value(items[0])
        if value and len(label) >= 8:
            candidates = items[1:]
    return sum(1 for item in candidates if _split_key_value(item)[1]) >= 3


def _looks_like_rich_cards(items: list[str]) -> bool:
    _lead, rest = _split_lead(items)
    pairs = [_split_key_value(item) for item in rest]
    bodies = [value for _label, value in pairs if value]
    if len(bodies) < 2:
        return False
    rich_bodies = [body for body in bodies if len(_split_card_points(body)) >= 2 or len(body) >= 54]
    return len(rich_bodies) >= 2


def _looks_like_process(title: str, items: list[str]) -> bool:
    joined = " ".join(items[:5])
    return (
        "怎么工作" in title
        or "如何工作" in title
        or "流程" in title
        or "Loop" in joined
        or joined.count("→") >= 2
        or joined.count("->") >= 2
    )


def _split_process_text(text: str) -> tuple[str, str]:
    title, body = _split_key_value(text)
    if body:
        return title, body
    arrow = "→" if "→" in text else "->" if "->" in text else ""
    if arrow:
        parts = [part.strip() for part in text.split(arrow) if part.strip()]
        if len(parts) >= 2:
            return parts[0], " → ".join(parts[1:])
    return _trim(text, 20), text


def _split_card_points(text: str) -> list[str]:
    text = _normalize_formula_text(text)
    if not text:
        return []
    parts = [part.strip(" 。；;") for part in re.split(r"\s*[；;]\s*", text) if part.strip(" 。；;")]
    if len(parts) == 1 and "。" in text:
        parts = [part.strip(" 。") for part in re.split(r"\s*。\s*", text) if part.strip(" 。")]
    return parts or [text]


def _title_size(title: str) -> float:
    if len(title) <= 18:
        return 20
    if len(title) <= 28:
        return 18
    return 16


def _trim(text: str, limit: int) -> str:
    text = text.strip()
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 1)].rstrip() + "…"
