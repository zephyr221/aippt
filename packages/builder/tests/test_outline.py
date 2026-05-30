import json
import os
import zipfile
from pathlib import Path

from aippt_builder.outline import outline_to_deck
from aippt_builder.render import build_pptx
from aippt_builder.validate import validate_deck
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx import Presentation
from pptx.util import Inches, Pt


def _sjtu_template_path() -> Path:
    candidates = [
        Path(value)
        for key in ("AIPPT_TEMPLATE_PPTX", "AIPPT_TEMPLATE_PPTX_PATH")
        if (value := os.getenv(key))
    ]
    candidates.extend(
        parent / "docs" / "SJTU PPT 模板" / "SJTU 模板.pptx"
        for parent in Path(__file__).resolve().parents
    )
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    raise AssertionError("SJTU template fixture is missing")


def test_outline_to_deck_creates_valid_buildable_ir() -> None:
    deck = outline_to_deck(
        """# AI PPT

## 背景

- 多用户登录
- 独立 PPT 工作区

## 生成链路

- Markdown outline
- Deck IR
- PPTX artifact
""",
        author="AIPPT",
    )

    assert deck.title == "AI PPT"
    assert deck.slides[0].layout == "cover"
    assert deck.slides[-1].layout == "thanks"
    assert validate_deck(deck) == []

    payload = json.loads(deck.model_dump_json())
    assert payload["slides"][1]["title"] == "背景"
    assert payload["slides"][2]["bullets"] == ["Markdown outline", "Deck IR", "PPTX artifact"]


def test_outline_promotes_cover_metadata_and_normalizes_page_headings() -> None:
    deck = outline_to_deck(
        """# AI × 计算材料科研

2026 年春，你需要知道的事
副标题：从对话到自主科研 Agent——理解当前 AI，用好当前 AI
SJTU · 计算材料课程组 组会分享
2026.03.27

## 第 1 页 · 封面

- 这行是旧封面备注，不应该生成正文页

## 第 2 页 · 开场：一个事实

- 三年前的今天，GPT-4 刚刚被看见
- 今天 Agent 已经能够拆解任务、调用工具、生成材料

## 第 3 页 · 我们为什么需要工作流

- 让生成过程可复现
- 让结果可以审阅、修复和下载

## 第 4 页 · 下一步

- 接入更强的 Agent 核心
- 加入模板和预览
""",
        author="AIPPT",
    )

    assert deck.slides[0].layout == "cover"
    assert "2026 年春" in deck.slides[0].subtitle
    assert "旧封面备注" not in deck.slides[0].subtitle
    assert [slide.title for slide in deck.slides[1:4]] == [
        "开场：一个事实",
        "我们为什么需要工作流",
        "下一步",
    ]
    assert all("第 " not in slide.title for slide in deck.slides[1:-1])
    assert all("封面" not in slide.title for slide in deck.slides[1:-1])
    assert validate_deck(deck) == []


def test_outline_limits_toc_to_valid_bullet_count() -> None:
    markdown = "# Long Deck\n\n" + "\n\n".join(
        f"## Section {idx}\n\n- Point {idx}" for idx in range(1, 8)
    )

    deck = outline_to_deck(markdown)

    assert deck.slides[1].layout == "toc"
    assert len(deck.slides[1].bullets) == 5
    assert validate_deck(deck) == []


def test_brief_prompt_expands_to_requested_intro_deck() -> None:
    deck = outline_to_deck("请制作五六页 PPT，关于机器学习的科普啊")

    assert deck.title == "机器学习科普"
    assert len(deck.slides) == 6
    assert [slide.title for slide in deck.slides] == [
        "机器学习科普",
        "为什么需要机器学习",
        "核心思想：从数据中学习规律",
        "最小例子：房价预测",
        "三类经典任务",
        "谢谢",
    ]
    body_text = "\n".join("\n".join(slide.bullets) for slide in deck.slides)
    assert "规则难以手写" in body_text
    assert "房价预测" in body_text
    assert "房价成交记录" in body_text
    assert "J(θ)=1/m" in body_text
    assert "预测 ŷ" in body_text
    assert "监督学习" in body_text
    assert "课堂练习" not in body_text
    assert len(deck.slides[1].bullets) == 4
    assert deck.slides[1].support == "规则系统的局限、数据中的规律和学习目标。"
    assert deck.slides[2].visual == "concept_diagram"
    assert deck.slides[3].visual == "example_walkthrough"
    assert all("：" in bullet for bullet in deck.slides[1].bullets[1:])
    assert all("；" in bullet for bullet in deck.slides[1].bullets[1:])
    assert validate_deck(deck) == []


def test_brief_machine_learning_intro_defaults_to_micro_lesson() -> None:
    deck = outline_to_deck("机器学习导论 PPT")

    assert deck.title == "机器学习导论"
    assert len(deck.slides) == 8
    assert [slide.title for slide in deck.slides[1:7]] == [
        "为什么需要机器学习",
        "核心思想：从数据中学习规律",
        "最小例子：房价预测",
        "三类经典任务",
        "训练流程与验证闭环",
        "常见误区与下一步",
    ]
    assert deck.slides[1].layout == "three_column"
    assert deck.slides[2].visual == "concept_diagram"
    assert deck.slides[3].visual == "example_walkthrough"
    assert deck.slides[6].visual == "summary"
    assert deck.slides[6].support == "误区、贯穿主线和学习路径共同收束。"
    body_text = "\n".join("\n".join(slide.bullets) for slide in deck.slides)
    assert "垃圾邮件" in body_text
    assert "贯穿主线" in body_text
    assert "课堂练习" not in body_text
    assert "训练集高分不等于真实可靠" in body_text
    assert validate_deck(deck) == []


def test_brief_prompt_handles_intro_deck_for_enterprise_audience() -> None:
    deck = outline_to_deck("为企业生成一份人工智能导论介绍 PPT")

    assert deck.title == "人工智能导论"
    assert len(deck.slides) == 8
    body_text = "\n".join("\n".join(slide.bullets) for slide in deck.slides)
    assert "学习 人工智能导论 之前" in body_text
    assert "小案例" in body_text
    assert "；" in body_text
    assert deck.slides[1].support == "用定义、例子和学习目标展开 人工智能导论。"
    assert deck.slides[-2].visual == "summary"
    assert validate_deck(deck) == []


def test_explicit_page_design_signals_create_structured_ir() -> None:
    deck = outline_to_deck(
        """# 结构化设计测试

## 第 1 页 · 封面

主标题：结构化设计测试
副标题：让 Hermes 选择页面组件

## 第 2 页 · 三类能力协同
版式：three_column
组件：rich_cards
支撑：规划日志、Deck IR 与预览 QA 会进入同一条生成链路。
三类能力共同决定 AI PPT 初稿是否可用。
- 规划能力：识别受众；建立主线；压缩标题
- 结构能力：选择版式；组织卡片；控制页面密度
- 验证能力：检查字数；渲染预览；反馈到下一轮
洞察：Hermes 负责设计选择，builder 负责稳定落地。

## 第 3 页 · 评估矩阵
版式：table
组件：table
- 维度：输入 / 输出 / 验证信号
- 内容规划：用户需求 / 页级大纲 / 标题是判断句
- 视觉设计：版式信号 / Deck IR / 组件节奏不重复
""",
    )

    card_slide = deck.slides[1]
    assert card_slide.layout == "three_column"
    assert card_slide.visual == "rich_cards"
    assert card_slide.proof == "规划日志、Deck IR 与预览 QA 会进入同一条生成链路。"
    assert card_slide.support == "规划日志、Deck IR 与预览 QA 会进入同一条生成链路。"
    assert len(card_slide.columns) == 3
    assert card_slide.columns[0].heading == "规划能力"
    assert card_slide.columns[0].bullets == ["识别受众", "建立主线", "压缩标题"]
    assert card_slide.insight == "Hermes 负责设计选择，builder 负责稳定落地。"

    table_slide = deck.slides[2]
    assert table_slide.layout == "table"
    assert table_slide.visual == "table"
    assert table_slide.table is not None
    assert table_slide.table.headers == ["维度", "输入", "输出", "验证信号"]
    assert table_slide.table.rows[0] == ["内容规划", "用户需求", "页级大纲", "标题是判断句"]
    assert validate_deck(deck) == []


def test_stat_callout_and_quote_block_render_as_editable_text(tmp_path) -> None:
    deck = outline_to_deck(
        """# 组件试验

## 第 1 页 · 封面
主标题：组件试验
副标题：指标和引用组件

## 第 2 页 · 质量提升先看三项指标
版式：one_column
组件：stat_callout
支撑：QA 报告、预览渲染与用户反馈会进入同一轮改写。
用指标约束生成质量，比只调 prompt 更可靠。
- 设计覆盖：100% / 内容页都要求组件信号
- 展开对象：每页 1 个 / 定义、步骤、案例、数据或公式
- 节奏重复：≤2 页 / 连续同版式会被 QA 提醒

## 第 3 页 · 设计原则
版式：one_column
组件：quote_block
支撑：来自当前 AIPPT 生成链路的稳定性约束。
模型做设计决策，builder 做稳定渲染。
- 原则：Hermes 负责版式和展开对象，Python 只执行白名单组件
- 落地：任何 repair 都必须先 validation，再 build，再 QA
""",
    )

    assert deck.slides[1].visual == "stat_callout"
    assert deck.slides[1].proof == "QA 报告、预览渲染与用户反馈会进入同一轮改写。"
    assert deck.slides[1].support == "QA 报告、预览渲染与用户反馈会进入同一轮改写。"
    assert deck.slides[2].visual == "quote_block"
    assert validate_deck(deck) == []

    output = build_pptx(deck, tmp_path / "components.pptx")
    with zipfile.ZipFile(output) as pptx:
        slide_xml = "\n".join(
            pptx.read(name).decode("utf-8")
            for name in pptx.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )

    assert "100%" in slide_xml
    assert "展开对象" in slide_xml
    assert "模型做设计决策" in slide_xml
    assert "白名单组件" in slide_xml


def test_report_components_render_as_editable_text(tmp_path) -> None:
    deck = outline_to_deck(
        """# 工作汇报组件

## 第 1 页 · 封面
主标题：工作汇报组件
副标题：吸收参考 PPT 的汇报型页面

## 第 2 页 · 三年工作总览
版式：one_column
组件：metric_strip
支撑：用关键数字和两条主线建立汇报全局。
- 系统数量：30+ / 独立开发系统与应用
- 课程支持：240+ / AI 课程深度支持
- 持续开发：16 / 个月持续 AI 开发
- 平台职责：平台 / AI 应用平台负责人

## 第 3 页 · AI 应用开发全景
版式：horizontal
组件：milestone_timeline
16 个月独立开发 30+ 系统，并全部稳定运行。
- 2024 秋：AI 翻译；AI 转录；本地 A100 部署
- 2024.11：AI 应用平台上线；招生 AI 审核；AI 组卷助手
- 2025.5：AI 修业导师；评教系统；研小知智能体
- 2025-2026：教学运行监控；查重系统；AI 知识库 / AIPPT

## 第 4 页 · 代表性项目
版式：horizontal
组件：project_showcase
国家级、省部级平台输出，并沉淀为自研产品。
- AI 组卷助手：上线国家智慧教育平台；中小学大规模使用
- 研小知智能体：为教育部学位中心建设；获得官方感谢信
- 研招自命题查重：保密环境独立开发；600+ 套试卷解析
- AI 知识库平台：高质量 RAG 引擎；教学与学科大模型底座

## 第 5 页 · AI 知识库平台
版式：horizontal
组件：media_explain
支撑：AI 知识库首页 / 笔记 / 检索界面截图。
对标腾讯 ima，沉淀教学与学科大模型的基础底座。
- 定位：自研高质量知识库平台，把文档做 RAG 处理，沉淀为可复用知识资产
- 学科大模型建设：向上支撑学科模型
- AI 教学材料与课程：服务教师制作和学生学习
- AI 知识库底座：提供 RAG 引擎和检索能力
""",
    )

    assert deck.slides[1].visual == "metric_strip"
    assert deck.slides[2].visual == "milestone_timeline"
    assert deck.slides[3].visual == "project_showcase"
    assert deck.slides[4].visual == "media_explain"
    assert validate_deck(deck) == []

    output = build_pptx(deck, tmp_path / "report-components.pptx")
    with zipfile.ZipFile(output) as pptx:
        slide_xml = "\n".join(
            pptx.read(name).decode("utf-8")
            for name in pptx.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )

    assert "30+" in slide_xml
    assert "2024 秋" in slide_xml
    assert "AI 组卷助手" in slide_xml
    assert "AI 知识库首页" in slide_xml
    assert "建议配图" in slide_xml


def test_formula_text_remains_editable_text_in_pptx(tmp_path) -> None:
    deck = outline_to_deck("请制作五六页 PPT，关于机器学习的科普啊")
    output = build_pptx(deck, tmp_path / "ml.pptx")

    with zipfile.ZipFile(output) as pptx:
        slide_xml = "\n".join(
            pptx.read(name).decode("utf-8")
            for name in pptx.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )

    assert "J(θ)=1/m" in slide_xml
    assert "Cambria Math" in slide_xml
    assert "预测 ŷ" in slide_xml
    assert "板书" in slide_xml
    assert "把误差变小" in slide_xml


def test_builder_uses_sjtu_template_when_configured(tmp_path, monkeypatch) -> None:
    template = _sjtu_template_path()
    monkeypatch.setenv("AIPPT_TEMPLATE_PPTX", str(template))
    deck = outline_to_deck("请制作五六页 PPT，关于机器学习的科普啊")
    output = build_pptx(deck, tmp_path / "template-backed.pptx")

    prs = Presentation(output)

    assert prs.slides[0].slide_layout.name == "封面-01"
    assert prs.slides[1].slide_layout.name == "常规样式（1）"
    assert prs.slides[-1].slide_layout.name == "封底01"
    assert 14 not in {
        shape.placeholder_format.idx
        for shape in prs.slides[1].shapes
        if shape.is_placeholder
    }
    cover_text = "\n".join(
        shape.text_frame.text
        for shape in prs.slides[0].shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    )
    assert "机器学习科普" in cover_text
    assert "6 页速览" in cover_text
    assert "一句话" not in cover_text
    assert len(prs.slides[1].shapes) >= 12
    first_body_text = "\n".join(
        shape.text_frame.text
        for shape in prs.slides[1].shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    )
    assert "▪" in first_body_text
    assert "规则方法" in first_body_text
    lead_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if shape.has_text_frame and "规则难以手写" in shape.text_frame.text
    )
    assert lead_shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert Inches(1.12) < lead_shape.top < Inches(1.2)
    point_sizes = [
        run.font.size
        for shape in prs.slides[1].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text.startswith("▪")
        for run in paragraph.runs
        if run.font.size is not None
    ]
    assert point_sizes
    assert min(point_sizes) >= Pt(11.5)
    process_title = next(
        shape
        for shape in prs.slides[3].shapes
        if shape.has_text_frame and shape.text_frame.text.strip() == "建立模型"
    )
    assert process_title.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
    process_arrow_types = []
    for shape in prs.slides[3].shapes:
        try:
            process_arrow_types.append(shape.auto_shape_type)
        except ValueError:
            continue
    assert MSO_SHAPE.RIGHT_ARROW not in process_arrow_types
    assert MSO_SHAPE.RIGHT_TRIANGLE not in process_arrow_types
    thanks = next(
        shape
        for shape in prs.slides[-1].shapes
        if shape.has_text_frame and shape.text_frame.text.strip() == "谢谢"
    )
    assert thanks.is_placeholder
    assert thanks.placeholder_format.idx == 11
    thanks_text = "\n".join(
        shape.text_frame.text
        for shape in prs.slides[-1].shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    )
    assert "单击此处添加文本" not in thanks_text
    assert Inches(0.85) < thanks.left < Inches(1.05)
    assert Inches(2.25) < thanks.top < Inches(2.45)


def test_explicit_page_outline_prefers_document_title_over_job_label() -> None:
    deck = outline_to_deck(
        """# 正式报告标题

## 第 1 页 · 封面

**正式报告标题**

## 第 2 页 · 内容

- 要点
""",
        title="AIPPT Demo",
    )

    assert deck.title == "正式报告标题"
    assert deck.slides[0].title == "正式报告标题"
    assert validate_deck(deck) == []


def test_explicit_page_outline_ignores_planning_notes_and_speaker_notes() -> None:
    deck = outline_to_deck(
        """# AI × 计算材料科研：2026 年春，你需要知道的事

> 面向 SJTU 计算材料课题组（DFT / MD）的组会分享
> 2026 年 3 月 27 日（周五）

---

## 我的整体思路

1. 这段是给讲者看的设计说明，不应该生成正文页。
2. 继续说明为什么要这样编排。

---

## 第 1 页 · 封面

# AI × 计算材料科研
**2026 年春，你需要知道的事**

副标题：从对话到自主科研 Agent——理解当前 AI，用好当前 AI

---

## 第 2 页 · 开场：一个事实

**三年前的今天，GPT-4 刚发布两周。**
**上周，一个 AI Agent 独自花了几天，从零写出了一个宇宙学 Boltzmann 求解器。**

- 2023.03.14 — GPT-4 发布
- 2026.03.23 — Claude 自主完成 Boltzmann solver
- 这类数值求解器，人类专家通常需要数月到数年

> 三年。从“能聊天”到“能独立做数值计算科研”。

**讲者备注**：
开场不讲道理，讲事实。这段不应进入 PPT 正文。

---

## 第 3 页 · AI 正在发生什么：一分钟理解核心脉络

| 阶段 | 环境 | 验证信号 |
|---|---|---|
| 早期 LLM | 无 | 人工标注 |
| Agent（当前） | 真实计算环境 | 任务是否完成 |

```text
这段 ASCII 图不应被当成 bullet。
```

> 对做 DFT/MD 的人来说，关键信息是精确验证。
""",
        author="AIPPT",
    )

    assert [slide.title for slide in deck.slides] == [
        "AI × 计算材料科研：2026 年春，你需要知道的事",
        "开场：一个事实",
        "AI 正在发生什么：一分钟理解核心脉络",
        "谢谢",
    ]
    body_text = "\n".join("\n".join(slide.bullets) for slide in deck.slides)
    assert "我的整体思路" not in body_text
    assert "讲者备注" not in body_text
    assert "ASCII 图" not in body_text
    assert "Agent（当前）" in body_text
    assert validate_deck(deck) == []
